# Copyright (c) 2026 nandan-d14. All rights reserved.
# Proprietary and non-commercial use only.

"""Modern dark widescreen deck renderer.

Runs standalone in the sandbox: python gen_pptx.py data.json out.pptx out.html
Paths may also come from DATA_PATH / OUT_PATH / HTML_PATH globals. Keep this
file free of ``from __future__`` imports and version-sensitive syntax so it
compiles on any Python 3 and survives any header prepended before it.
"""

import html as html_lib
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BG = RGBColor(0x0B, 0x0B, 0x0E)
SURFACE = RGBColor(0x16, 0x16, 0x1C)
TEXT = RGBColor(0xF5, 0xF5, 0xF4)
MUTED = RGBColor(0xA8, 0xA2, 0x9E)
ACCENT = RGBColor(0x2D, 0xD4, 0xBF)
RULE = RGBColor(0x29, 0x25, 0x24)

FONT = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.72)
MARGIN_Y = Inches(0.48)


def _blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if str(getattr(layout, "name", "")).strip().lower() == "blank":
            return layout
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]


def _no_line(shape) -> None:
    try:
        shape.line.fill.background()
    except Exception:
        pass


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    _no_line(shape)


def _set_run(paragraph, text: str, *, size: int, bold: bool = False, color=TEXT, align=None) -> None:
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.font.name = FONT
    if align is not None:
        paragraph.alignment = align


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _kicker_and_title(slide, item: dict, *, title_size: int = 36, top=None) -> int:
    top = MARGIN_Y if top is None else top
    kicker = str(item.get("kicker") or "").strip()
    title = str(item.get("title") or "").strip()
    subtitle = str(item.get("subtitle") or "").strip()
    y = top
    if kicker:
        tf = _textbox(slide, MARGIN_X, y, Inches(11.8), Inches(0.32))
        _set_run(tf.paragraphs[0], kicker.upper(), size=12, bold=True, color=ACCENT)
        y += Inches(0.34)
    title_h = Inches(1.15) if len(title) < 48 else Inches(1.55)
    tf = _textbox(slide, MARGIN_X, y, Inches(11.9), title_h)
    _set_run(tf.paragraphs[0], title, size=title_size, bold=True, color=TEXT)
    y += title_h
    if subtitle:
        tf = _textbox(slide, MARGIN_X, y, Inches(11.6), Inches(0.55))
        _set_run(tf.paragraphs[0], subtitle, size=16, color=MUTED)
        y += Inches(0.5)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X, y + Inches(0.06), Inches(1.35), Pt(3))
    _fill(rule, ACCENT)
    return y + Inches(0.28)


def _bullets(slide, lines: list, left, top, width, height, *, size: int = 20) -> None:
    if not lines:
        return
    tf = _textbox(slide, left, top, width, height)
    if hasattr(tf, "clear"):
        tf.clear()
    for index, line in enumerate(lines[:8]):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.space_after = Pt(10)
        paragraph.level = 0
        _set_run(paragraph, f"▸  {line}", size=size, color=TEXT)


def _footer(slide, footnote: str, page: int, total: int) -> None:
    label = footnote or ""
    tf = _textbox(slide, MARGIN_X, Inches(7.12), Inches(9.4), Inches(0.28))
    _set_run(tf.paragraphs[0], label, size=11, color=MUTED)
    tf = _textbox(slide, Inches(10.6), Inches(7.12), Inches(2.0), Inches(0.28))
    _set_run(tf.paragraphs[0], f"{page:02d}  /  {total:02d}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def _chrome(slide) -> None:
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG
    except Exception:
        backdrop = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        _fill(backdrop, BG)
        sp_tree = slide.shapes._spTree
        element = backdrop._element
        sp_tree.remove(element)
        sp_tree.insert(2, element)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), SLIDE_H)
    _fill(bar, ACCENT)


def _card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill(shape, SURFACE)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    try:
        shape.line.color.rgb = RULE
        shape.line.width = Pt(0.75)
    except Exception:
        pass
    return shape


def render_title(slide, item: dict) -> None:
    kicker = str(item.get("kicker") or "").strip()
    if kicker:
        tf = _textbox(slide, MARGIN_X, Inches(1.85), Inches(11.8), Inches(0.36))
        _set_run(tf.paragraphs[0], kicker.upper(), size=13, bold=True, color=ACCENT)
    tf = _textbox(slide, MARGIN_X, Inches(2.28), Inches(12.0), Inches(2.1))
    _set_run(tf.paragraphs[0], str(item.get("title") or ""), size=48, bold=True, color=TEXT)
    subtitle = str(item.get("subtitle") or "").strip()
    if subtitle:
        tf = _textbox(slide, MARGIN_X, Inches(4.5), Inches(11.2), Inches(0.9))
        _set_run(tf.paragraphs[0], subtitle, size=20, color=MUTED)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(4.32), Inches(1.6), Pt(4))
    _fill(rule, ACCENT)


def render_section(slide, item: dict, index: int) -> None:
    tf = _textbox(slide, MARGIN_X, Inches(2.05), Inches(12), Inches(1.2))
    _set_run(tf.paragraphs[0], f"{index:02d}", size=54, bold=True, color=ACCENT)
    tf = _textbox(slide, MARGIN_X, Inches(3.2), Inches(12), Inches(1.6))
    _set_run(tf.paragraphs[0], str(item.get("title") or ""), size=40, bold=True, color=TEXT)
    subtitle = str(item.get("subtitle") or "").strip()
    if subtitle:
        tf = _textbox(slide, MARGIN_X, Inches(4.9), Inches(11), Inches(0.8))
        _set_run(tf.paragraphs[0], subtitle, size=18, color=MUTED)


def render_content(slide, item: dict) -> None:
    y = _kicker_and_title(slide, item, title_size=32)
    bullets = [str(b) for b in (item.get("bullets") or []) if str(b).strip()]
    size = 18 if len(bullets) > 5 else 20
    _bullets(slide, bullets, MARGIN_X, y + Inches(0.12), Inches(11.8), Inches(4.4), size=size)


def render_split(slide, item: dict) -> None:
    y = _kicker_and_title(slide, item, title_size=30)
    left = [str(b) for b in (item.get("left") or item.get("bullets") or []) if str(b).strip()]
    right = [str(b) for b in (item.get("right") or []) if str(b).strip()]
    gap = Inches(0.28)
    width = Inches(5.8)
    height = Inches(4.15)
    top = y + Inches(0.18)
    _card(slide, MARGIN_X, top, width, height)
    _card(slide, MARGIN_X + width + gap, top, width, height)
    _bullets(slide, left[:6], MARGIN_X + Inches(0.28), top + Inches(0.28), width - Inches(0.5), height - Inches(0.45), size=16)
    _bullets(slide, right[:6], MARGIN_X + width + gap + Inches(0.28), top + Inches(0.28), width - Inches(0.5), height - Inches(0.45), size=16)


def render_stats(slide, item: dict) -> None:
    y = _kicker_and_title(slide, item, title_size=30)
    stats = list(item.get("stats") or [])[:4]
    if not stats:
        return
    count = max(len(stats), 1)
    gap = Inches(0.22)
    width = (Inches(11.9) - gap * (count - 1)) / count
    top = y + Inches(0.35)
    height = Inches(3.4)
    for index, stat in enumerate(stats):
        left = MARGIN_X + (width + gap) * index
        _card(slide, left, top, width, height)
        tf = _textbox(slide, left + Inches(0.28), top + Inches(0.7), width - Inches(0.5), Inches(1.3))
        _set_run(tf.paragraphs[0], str(stat.get("value") or ""), size=36, bold=True, color=ACCENT)
        tf = _textbox(slide, left + Inches(0.28), top + Inches(2.1), width - Inches(0.5), Inches(0.9))
        _set_run(tf.paragraphs[0], str(stat.get("label") or ""), size=15, color=MUTED)


def render_quote(slide, item: dict) -> None:
    tf = _textbox(slide, MARGIN_X, Inches(1.35), Inches(12), Inches(0.7))
    _set_run(tf.paragraphs[0], "“", size=72, bold=True, color=ACCENT)
    quote = str(item.get("quote") or item.get("title") or "")
    tf = _textbox(slide, MARGIN_X, Inches(2.3), Inches(11.8), Inches(2.8))
    _set_run(tf.paragraphs[0], quote, size=28, bold=True, color=TEXT)
    attribution = str(item.get("attribution") or "").strip()
    if attribution:
        tf = _textbox(slide, MARGIN_X, Inches(5.35), Inches(11), Inches(0.5))
        _set_run(tf.paragraphs[0], attribution, size=16, color=MUTED)


def render_closing(slide, item: dict) -> None:
    kicker = str(item.get("kicker") or "Next").strip()
    tf = _textbox(slide, MARGIN_X, Inches(1.7), Inches(12), Inches(0.36))
    _set_run(tf.paragraphs[0], kicker.upper(), size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf = _textbox(slide, MARGIN_X, Inches(2.15), Inches(12), Inches(1.5))
    _set_run(tf.paragraphs[0], str(item.get("title") or "Thank you"), size=44, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    bullets = [str(b) for b in (item.get("bullets") or []) if str(b).strip()]
    if bullets:
        _bullets(slide, bullets[:5], Inches(3.4), Inches(3.9), Inches(6.5), Inches(2.6), size=18)


def render_slide(slide, item: dict, index: int, total: int) -> None:
    _chrome(slide)
    layout = str(item.get("layout") or "content")
    if layout == "title":
        render_title(slide, item)
    elif layout == "section":
        render_section(slide, item, index)
    elif layout == "split":
        render_split(slide, item)
    elif layout == "stats":
        render_stats(slide, item)
    elif layout == "quote":
        render_quote(slide, item)
    elif layout == "closing":
        render_closing(slide, item)
    else:
        render_content(slide, item)
    if layout != "title":
        _footer(slide, str(item.get("footnote") or ""), index + 1, total)


def _escape(value: Any) -> str:
    return html_lib.escape(str(value or ""))


def render_html(data: dict) -> str:
    slides = data.get("slides") or []
    sections = []
    total = len(slides)
    for index, item in enumerate(slides):
        layout = _escape(item.get("layout") or "content")
        kicker = _escape(item.get("kicker") or "")
        title = _escape(item.get("title") or "")
        subtitle = _escape(item.get("subtitle") or "")
        quote = _escape(item.get("quote") or "")
        attribution = _escape(item.get("attribution") or "")
        footnote = _escape(item.get("footnote") or "")
        bullets = "".join(f"<li>{_escape(b)}</li>" for b in (item.get("bullets") or []) if str(b).strip())
        left = "".join(f"<li>{_escape(b)}</li>" for b in (item.get("left") or []) if str(b).strip())
        right = "".join(f"<li>{_escape(b)}</li>" for b in (item.get("right") or []) if str(b).strip())
        stats = "".join(
            f'<div class="stat"><div class="stat-value">{_escape(s.get("value"))}</div>'
            f'<div class="stat-label">{_escape(s.get("label"))}</div></div>'
            for s in (item.get("stats") or [])
        )
        kicker_html = f'<p class="kicker">{kicker}</p>' if kicker else ""
        subtitle_html = f'<p class="subtitle">{subtitle}</p>' if subtitle else ""
        bullets_html = f"<ul>{bullets}</ul>" if bullets else ""
        split_html = ""
        if left or right:
            split_html = (
                f'<div class="split"><div class="card"><ul>{left}</ul></div>'
                f'<div class="card"><ul>{right}</ul></div></div>'
            )
        stats_html = f'<div class="stats">{stats}</div>' if stats else ""
        quote_html = f'<blockquote>{quote}</blockquote>' if quote else ""
        cite_html = f'<p class="cite">{attribution}</p>' if attribution else ""
        sections.append(
            f'<section class="slide layout-{layout}">'
            f'<div class="inner">{kicker_html}<h1>{title}</h1>{subtitle_html}'
            f"{quote_html}{cite_html}{bullets_html}{split_html}{stats_html}</div>"
            f'<div class="foot"><span>{footnote}</span><span>{index + 1:02d} / {total:02d}</span></div>'
            "</section>"
        )
    deck_title = _escape(data.get("title") or "Slides")
    return f"""<!DOCTYPE html>
<html lang="en"><head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{deck_title}</title>
<style>
:root {{
  --bg: #0B0B0E; --surface: #16161C; --text: #F5F5F4; --muted: #A8A29E;
  --accent: #2DD4BF; --rule: #292524;
}}
* {{ box-sizing: border-box; }}
html, body {{ margin: 0; background: var(--bg); color: var(--text);
  font-family: Calibri, "Segoe UI", Inter, system-ui, sans-serif; }}
.slide {{
  min-height: 100vh; padding: 56px 72px 48px; position: relative;
  border-bottom: 1px solid var(--rule); display: flex; flex-direction: column;
  justify-content: center; border-left: 6px solid var(--accent);
}}
.inner {{ max-width: 1100px; }}
.kicker {{ color: var(--accent); letter-spacing: 0.22em; text-transform: uppercase;
  font-size: 0.78rem; font-weight: 700; margin: 0 0 12px; }}
h1 {{ font-size: clamp(2rem, 4.4vw, 3.4rem); letter-spacing: -0.035em; line-height: 1.08;
  margin: 0 0 16px; font-weight: 700; }}
.subtitle {{ color: var(--muted); font-size: 1.2rem; margin: 0 0 28px; max-width: 46rem; }}
ul {{ margin: 8px 0 0; padding: 0; list-style: none; font-size: 1.22rem; line-height: 1.45; }}
li {{ margin: 0.55rem 0; padding-left: 1.15rem; position: relative; }}
li::before {{ content: "▸"; color: var(--accent); position: absolute; left: 0; }}
.split {{ display: grid; grid-template-columns: 1fr 1fr; gap: 18px; margin-top: 12px; }}
.card {{ background: var(--surface); border: 1px solid var(--rule); border-radius: 18px; padding: 22px 24px; }}
.stats {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(180px, 1fr)); gap: 16px; margin-top: 10px; }}
.stat {{ background: var(--surface); border: 1px solid var(--rule); border-radius: 18px; padding: 28px 24px; }}
.stat-value {{ color: var(--accent); font-size: 2.4rem; font-weight: 700; letter-spacing: -0.03em; }}
.stat-label {{ color: var(--muted); margin-top: 8px; }}
blockquote {{ font-size: 2rem; font-weight: 600; letter-spacing: -0.03em; line-height: 1.25; margin: 0; }}
.cite {{ color: var(--muted); margin-top: 18px; }}
.layout-title h1 {{ font-size: clamp(2.6rem, 6vw, 4.2rem); }}
.layout-title, .layout-section, .layout-closing {{ justify-content: center; }}
.layout-closing h1, .layout-closing .kicker {{ text-align: center; }}
.layout-closing ul {{ max-width: 36rem; margin-left: auto; margin-right: auto; }}
.foot {{ position: absolute; left: 72px; right: 72px; bottom: 28px; display: flex;
  justify-content: space-between; color: var(--muted); font-size: 0.85rem; letter-spacing: 0.04em; }}
</style></head>
<body>{''.join(sections)}</body></html>
"""


def build_deck(data: dict, out_path: str, html_path: str) -> dict:
    slides = data.get("slides") or []
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    layout = _blank_layout(prs)
    total = len(slides)
    for index, item in enumerate(slides):
        slide = prs.slides.add_slide(layout)
        render_slide(slide, item, index, total)
    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)
    html_doc = render_html(data)
    os.makedirs(os.path.dirname(html_path) or ".", exist_ok=True)
    with open(html_path, "w", encoding="utf-8") as handle:
        handle.write(html_doc)
    return {
        "status": "success",
        "path": out_path,
        "html_path": html_path,
        "size": os.path.getsize(out_path),
        "html_size": os.path.getsize(html_path),
        "slide_count": total,
        "theme": "modern-dark",
    }


def main() -> None:
    data_path = globals().get("DATA_PATH") or sys.argv[1]
    out_path = globals().get("OUT_PATH") or sys.argv[2]
    html_path = globals().get("HTML_PATH") or sys.argv[3]
    with open(data_path, encoding="utf-8") as handle:
        data = json.load(handle)
    print(json.dumps(build_deck(data, out_path, html_path)))


if __name__ == "__main__":
    main()
