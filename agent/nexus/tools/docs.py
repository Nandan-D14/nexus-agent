# Copyright (c) 2026 Agentic Company. All rights reserved.
# Proprietary and non-commercial use only.

"""Document generation and artifact management tools."""

from __future__ import annotations

import asyncio
import base64
import inspect
import json
import logging
import os
import re
import shlex
import textwrap
from typing import Any

from nexus.document_design import PDF_REPORT_CSS, normalize_slides, pptx_generator_source
from nexus.tools.base import normalized_tool, tool_error, tool_success
from nexus.tools._context import (
    get_artifact_callback,
    get_history_repository,
    get_run_id,
    get_sandbox,
    get_session_id,
    get_workspace_path,
)
from nexus.tools.sandbox_events import emit_sandbox_event
from nexus.storage import artifact_storage_metadata, upload_artifact_async

logger = logging.getLogger(__name__)

_HTML_DATA_URI_LIMIT_BYTES = 500_000

# #region agent log
def _dbg_pptx_log(hypothesis_id: str, location: str, message: str, data: dict[str, Any]) -> None:
    payload = {
        "sessionId": "993e46",
        "runId": "post-fix",
        "hypothesisId": hypothesis_id,
        "location": location,
        "message": message,
        "data": data,
        "timestamp": int(__import__("time").time() * 1000),
    }
    line = json.dumps(payload)
    try:
        with open(
            r"c:\Users\nanda\OneDrive\Desktop\co-computer\debug-993e46.log",
            "a",
            encoding="utf-8",
        ) as handle:
            handle.write(line + "\n")
    except Exception:
        pass
    try:
        import urllib.request

        req = urllib.request.Request(
            "http://127.0.0.1:7421/ingest/08b059be-2c03-45ae-97a1-bb3c6f862ec1",
            data=line.encode("utf-8"),
            headers={
                "Content-Type": "application/json",
                "X-Debug-Session-Id": "993e46",
            },
            method="POST",
        )
        urllib.request.urlopen(req, timeout=1).read()
    except Exception:
        pass
# #endregion


# Prepended to sandbox generator scripts so missing packages self-heal
# when boot-time provisioning was skipped or failed silently.
_SANDBOX_DEPS_BOOTSTRAP = textwrap.dedent("""\
    import importlib, subprocess, sys
    for pkg, mod in [("python-docx", "docx"), ("openpyxl", "openpyxl"),
                     ("python-pptx", "pptx"), ("weasyprint", "weasyprint"),
                     ("fpdf2", "fpdf"), ("markdown2", "markdown2")]:
        try:
            importlib.import_module(mod)
        except ImportError:
            subprocess.run([sys.executable, "-m", "pip", "install", pkg],
                           check=True, capture_output=True, timeout=240)
""")


def _safe_html_filename(value: str | None, *, fallback: str = "artifact.html") -> str:
    cleaned = (value or fallback).strip().replace("\\", "/").rsplit("/", 1)[-1]
    cleaned = re.sub(r"[^A-Za-z0-9._ -]+", "-", cleaned).strip(" .-")
    if not cleaned:
        cleaned = fallback
    if not cleaned.lower().endswith(".html"):
        cleaned = f"{cleaned}.html"
    return cleaned


def _ensure_full_html_document(title: str, html_content: str) -> str:
    body = html_content.strip()
    if re.search(r"<!doctype\s+html|<html[\s>]", body, flags=re.IGNORECASE):
        return body
    safe_title = title.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    return (
        "<!doctype html>\n"
        '<html lang="en">\n'
        "<head>\n"
        '  <meta charset="utf-8">\n'
        '  <meta name="viewport" content="width=device-width, initial-scale=1">\n'
        f"  <title>{safe_title}</title>\n"
        "</head>\n"
        "<body>\n"
        f"{body}\n"
        "</body>\n"
        "</html>\n"
    )


def _html_preview_text(html_content: str) -> str:
    compact = re.sub(r"<(script|style)[\s\S]*?</\1>", " ", html_content, flags=re.IGNORECASE)
    compact = re.sub(r"<[^>]+>", " ", compact)
    compact = " ".join(compact.split())
    return compact[:240] if compact else "Interactive HTML artifact ready."


def _normalize_slides(slides: list[Any] | None, *, deck_title: str = "") -> list[dict[str, Any]]:
    return normalize_slides(slides, deck_title=deck_title)


def _artifact_payload(artifact) -> dict[str, Any]:
    return {
        "artifact_id": artifact.artifact_id,
        "run_id": artifact.run_id,
        "session_id": artifact.session_id,
        "task_id": getattr(artifact, "task_id", None),
        "kind": artifact.kind,
        "title": artifact.title,
        "preview": artifact.preview,
        "created_at": artifact.created_at.isoformat() if getattr(artifact, "created_at", None) else None,
        "source_step_id": artifact.source_step_id,
        "path": artifact.path,
        "url": artifact.url,
        "metadata": artifact.metadata or {},
    }


async def _notify_artifact_created(artifact) -> None:
    callback = get_artifact_callback()
    if callback is None:
        return
    result = callback(_artifact_payload(artifact))
    if inspect.isawaitable(result):
        await result


def _resolve_workspace_or_absolute_path(path: str) -> str:
    cleaned = (path or "").strip()
    if not cleaned:
        raise ValueError("path is required")
    if cleaned.startswith("/"):
        return cleaned
    return f"{get_workspace_path().rstrip('/')}/{cleaned.lstrip('/')}"


@normalized_tool(needs_sandbox=True)
async def extract_pdf_text(path: str, max_chars: int = 12000) -> dict[str, Any]:
    """Extract readable text from a PDF upload or a base64 text file containing a PDF."""
    sandbox = get_sandbox()
    source_path = _resolve_workspace_or_absolute_path(path)
    max_chars = max(1000, min(int(max_chars or 12000), 30000))
    output_dir = f"{get_workspace_path().rstrip('/')}/sources/pdf_text"
    sandbox.ensure_directory(output_dir)
    base_name = os.path.basename(source_path).rsplit(".", 1)[0] or "document"
    output_path = f"{output_dir}/{base_name}.txt"

    script = f"""
import base64
import json
import os
import subprocess
import sys
import tempfile
from pathlib import Path

source = Path({source_path!r})
out_path = Path({output_path!r})
max_chars = {max_chars!r}

def maybe_decode_pdf(path: Path) -> Path:
    raw = path.read_bytes()
    if raw.startswith(b"%PDF"):
        return path
    text = raw.decode("ascii", errors="ignore").strip()
    if text.startswith("JVBER"):
        decoded = base64.b64decode(text)
        tmp = Path(tempfile.gettempdir()) / (path.stem + "_decoded.pdf")
        tmp.write_bytes(decoded)
        return tmp
    return path

def extract_with_python(path: Path) -> str:
    errors = []
    for module_name in ("pypdf", "PyPDF2"):
        try:
            module = __import__(module_name)
            reader = module.PdfReader(str(path))
            pages = []
            for i, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"--- Page {{i}} ---\\n{{text.strip()}}")
            return "\\n\\n".join(pages)
        except Exception as exc:
            errors.append(f"{{module_name}}: {{exc}}")
    raise RuntimeError("; ".join(errors))

def extract_with_pdftotext(path: Path) -> str:
    proc = subprocess.run(
        ["pdftotext", "-layout", str(path), "-"],
        text=True,
        capture_output=True,
        timeout=60,
    )
    if proc.returncode != 0:
        raise RuntimeError(proc.stderr.strip() or "pdftotext failed")
    return proc.stdout

try:
    pdf_path = maybe_decode_pdf(source)
    text = ""
    errors = []
    for extractor in (extract_with_python, extract_with_pdftotext):
        try:
            text = extractor(pdf_path)
            if text.strip():
                break
        except Exception as exc:
            errors.append(str(exc))
    if not text.strip():
        raise RuntimeError("; ".join(errors) or "No text could be extracted")
    full_text = text.strip()
    out_path.write_text(full_text, encoding="utf-8")
    print(json.dumps({{
        "status": "success",
        "source_path": str(source),
        "text_path": str(out_path),
        "char_count": len(full_text),
        "text_excerpt": full_text[:max_chars],
        "truncated": len(full_text) > max_chars,
    }}))
except Exception as exc:
    print(json.dumps({{
        "status": "error",
        "message": str(exc),
        "source_path": str(source),
    }}))
    sys.exit(1)
"""

    script_path = f"/tmp/extract_pdf_text_{get_run_id()}.py"
    sandbox.write_text_file(script_path, script)
    result = sandbox.run_command(f"python3 {shlex.quote(script_path)}", timeout=90)
    stdout = str(result.get("stdout") or "").strip()
    decoded_path = f"/tmp/{base_name}_decoded.pdf"
    sandbox.run_command(
        f"rm -f {shlex.quote(script_path)} {shlex.quote(decoded_path)}",
        timeout=10,
    )
    if result.get("exit_code") != 0:
        return {
            "status": "error",
            "summary": f"PDF text extraction failed: {result.get('stderr') or stdout}",
            "detail": result,
        }
    try:
        import json

        payload = json.loads(stdout)
    except Exception:
        return {
            "status": "error",
            "summary": "PDF text extraction returned invalid output.",
            "detail": result,
        }
    if payload.get("status") != "success":
        return {
            "status": "error",
            "summary": payload.get("message") or "PDF text extraction failed.",
            "detail": payload,
        }
    return {
        "status": "success",
        "summary": f"Extracted PDF text to {payload.get('text_path')}",
        "detail": payload,
        "metadata": {
            "source_path": payload.get("source_path"),
            "text_path": payload.get("text_path"),
            "char_count": payload.get("char_count"),
            "truncated": payload.get("truncated"),
        },
    }

_OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx", ".csv"}


@normalized_tool(needs_sandbox=True)
async def extract_document_text(path: str, max_chars: int = 12000) -> dict[str, Any]:
    """Extract readable text from a .docx, .xlsx, .pptx, or .csv file.

    Use this on Office uploads before trying to reason about them. The full text
    is also written under sources/document_text/ so search_sources can find it.

    Args:
        path: Workspace-relative or absolute path to the document.
        max_chars: Maximum characters to return inline (1000-30000).

    Returns:
        NormalizedToolResult with a text excerpt and the extracted-text path.
    """
    sandbox = get_sandbox()
    source_path = _resolve_workspace_or_absolute_path(path)
    extension = os.path.splitext(source_path)[1].lower()
    if extension not in _OFFICE_EXTENSIONS:
        return tool_error(
            f"{extension or 'This file type'} is not supported by "
            "extract_document_text. Use extract_pdf_text for PDFs or "
            "read_workspace_file for plain text.",
            error_code="UNSUPPORTED_FORMAT",
            suggested_alternatives=["extract_pdf_text", "read_workspace_file"],
        )

    max_chars = max(1000, min(int(max_chars or 12000), 30000))
    output_dir = f"{get_workspace_path().rstrip('/')}/sources/document_text"
    sandbox.ensure_directory(output_dir)
    base_name = os.path.basename(source_path).rsplit(".", 1)[0] or "document"
    output_path = f"{output_dir}/{base_name}.txt"

    script = f"""
import json
import sys
from pathlib import Path

source = Path({source_path!r})
out_path = Path({output_path!r})
max_chars = {max_chars!r}
suffix = source.suffix.lower()

def from_docx(path):
    import docx
    document = docx.Document(str(path))
    blocks = [p.text for p in document.paragraphs if p.text.strip()]
    for table_index, table in enumerate(document.tables, start=1):
        blocks.append(f"--- Table {{table_index}} ---")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                blocks.append(" | ".join(cells))
    return "\\n".join(blocks)

def from_xlsx(path):
    import openpyxl
    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    blocks = []
    for sheet in workbook.worksheets:
        blocks.append(f"--- Sheet: {{sheet.title}} ---")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if value is None else str(value) for value in row]
            if any(cell.strip() for cell in cells):
                blocks.append(" | ".join(cells))
    workbook.close()
    return "\\n".join(blocks)

def from_pptx(path):
    from pptx import Presentation
    presentation = Presentation(str(path))
    blocks = []
    for index, slide in enumerate(presentation.slides, start=1):
        blocks.append(f"--- Slide {{index}} ---")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                blocks.append(text.strip())
    return "\\n".join(blocks)

def from_csv(path):
    return path.read_text(encoding="utf-8", errors="replace")

extractors = {{
    ".docx": from_docx,
    ".xlsx": from_xlsx,
    ".pptx": from_pptx,
    ".csv": from_csv,
}}

try:
    text = (extractors[suffix](source) or "").strip()
    if not text:
        raise RuntimeError("The document contains no extractable text.")
    out_path.write_text(text, encoding="utf-8")
    print(json.dumps({{
        "status": "success",
        "source_path": str(source),
        "text_path": str(out_path),
        "format": suffix.lstrip("."),
        "char_count": len(text),
        "text_excerpt": text[:max_chars],
        "truncated": len(text) > max_chars,
    }}))
except Exception as exc:
    print(json.dumps({{
        "status": "error",
        "message": str(exc),
        "source_path": str(source),
    }}))
    sys.exit(1)
"""

    script_path = f"/tmp/extract_document_text_{get_run_id()}.py"
    sandbox.write_text_file(script_path, script)
    result = sandbox.run_command(f"python3 {shlex.quote(script_path)}", timeout=120)
    stdout = str(result.get("stdout") or "").strip()
    sandbox.run_command(f"rm -f {shlex.quote(script_path)}", timeout=10)

    if result.get("exit_code") != 0:
        return tool_error(
            f"Document text extraction failed: {result.get('stderr') or stdout}",
            error_code="EXTRACTION_FAILED",
            retryable=False,
        )
    try:
        payload = json.loads(stdout)
    except Exception:
        return tool_error(
            "Document text extraction returned invalid output.",
            error_code="EXTRACTION_FAILED",
        )
    if payload.get("status") != "success":
        return tool_error(
            str(payload.get("message") or "Document text extraction failed."),
            error_code="EXTRACTION_FAILED",
        )
    return tool_success(
        f"Extracted {payload.get('format')} text to {payload.get('text_path')}",
        detail=payload,
        source_path=payload.get("source_path"),
        text_path=payload.get("text_path"),
        char_count=payload.get("char_count"),
        truncated=payload.get("truncated"),
    )


@normalized_tool(needs_sandbox=True)
async def generate_pdf_report(
    title: str,
    markdown_content: str,
    filename: str | None = None
) -> dict[str, Any]:
    """
    Generate a professional PDF report from Markdown content.

    Use structured markdown: H2 sections, short paragraphs, tables, and
    bullet lists. Do not wrap the whole report in a single code block.

    Args:
        title: The title of the report (appears on a dark cover band).
        markdown_content: The body of the report in Markdown format.
        filename: Optional desired filename (e.g., 'analysis.pdf').
    """
    sandbox = get_sandbox()
    session_id = get_session_id()
    run_id = get_run_id()
    history_repo = get_history_repository()

    if not filename:
        filename = f"report_{run_id[:8]}.pdf"
    if not filename.endswith(".pdf"):
        filename += ".pdf"
    filename = os.path.basename(filename)
    output_relative_path = f"outputs/{filename}"
    output_path = f"{get_workspace_path().rstrip('/')}/{output_relative_path}"

    # Write the markdown content to a temp file to avoid shell quoting issues
    md_temp = f"/tmp/_pdf_md_{run_id}.md"
    sandbox.write_text_file(md_temp, markdown_content)

    # Write the PDF generation script
    pdf_script = _SANDBOX_DEPS_BOOTSTRAP + textwrap.dedent(f"""\
import json, os, sys, pathlib

title = {repr(title)}
md_path = {repr(md_temp)}
out_path = {repr(output_path)}
CSS_TEMPLATE = {json.dumps(PDF_REPORT_CSS)}

os.makedirs(os.path.dirname(out_path), exist_ok=True)
md_text = pathlib.Path(md_path).read_text(encoding="utf-8")

def try_weasyprint(md_text, title, out_path):
    import markdown2
    from weasyprint import HTML, CSS
    html_body = markdown2.markdown(
        md_text,
        extras=["fenced-code-blocks", "tables", "break-on-newline",
                "header-ids", "strike", "task_list"]
    )
    full_html = f\"\"\"<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>{{title}}</title></head>
<body>
<header class="cover"><p class="kicker">Report</p><h1>{{title}}</h1></header>
{{html_body}}
</body></html>\"\"\"
    HTML(string=full_html).write_pdf(out_path, stylesheets=[CSS(string=CSS_TEMPLATE)])
    return True

def try_fpdf(md_text, title, out_path):
    from fpdf import FPDF
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, title.encode("latin-1", errors="replace").decode("latin-1"), ln=True, align="C")
    pdf.ln(8)
    pdf.set_font("Helvetica", size=11)
    for line in md_text.splitlines():
        safe = line.encode("latin-1", errors="replace").decode("latin-1")
        if line.startswith("## "):
            pdf.ln(4)
            pdf.set_font("Helvetica", "B", 14)
            pdf.multi_cell(0, 7, safe[3:])
            pdf.set_font("Helvetica", size=11)
        elif line.startswith("# "):
            pdf.ln(4)
            pdf.set_font("Helvetica", "B", 16)
            pdf.multi_cell(0, 8, safe[2:])
            pdf.set_font("Helvetica", size=11)
        elif line.strip():
            pdf.multi_cell(0, 6, safe)
        else:
            pdf.ln(3)
    pdf.output(out_path)
    return True

errors = []
for gen_fn in (try_weasyprint, try_fpdf):
    try:
        gen_fn(md_text, title, out_path)
        size = os.path.getsize(out_path)
        print(json.dumps({{"status": "success", "path": out_path, "size": size,
                           "engine": gen_fn.__name__}}))
        sys.exit(0)
    except Exception as exc:
        errors.append(f"{{gen_fn.__name__}}: {{exc}}")

print(json.dumps({{"status": "error", "message": "; ".join(errors)}}))
sys.exit(1)
""")

    script_path = f"/tmp/gen_pdf_{run_id}.py"
    sandbox.write_text_file(script_path, pdf_script)

    res = sandbox.run_command(f"python3 {script_path}", timeout=120)
    sandbox.run_command(f"rm -f {shlex.quote(script_path)} {shlex.quote(md_temp)}", timeout=10)

    if res.get("exit_code") != 0:
        return {
            "status": "error",
            "summary": f"PDF generation failed: {res.get('stderr') or res.get('stdout')}",
            "detail": res,
        }

    try:
        import json as _json
        payload = _json.loads(str(res.get("stdout") or "").strip())
    except Exception:
        payload = {}

    if payload.get("status") != "success":
        return {
            "status": "error",
            "summary": payload.get("message") or "PDF generation returned unexpected output.",
            "detail": res,
        }

    # Upload to GCS
    try:
        content = sandbox.read_binary_file(output_path)
        gcs_url = await upload_artifact_async(
            session_id=session_id,
            run_id=run_id,
            relative_path=output_relative_path,
            content=content,
        )

        artifact = await history_repo.create_artifact(
            session_id=session_id,
            run_id=run_id,
            kind="pdf_report",
            title=title or filename,
            preview=f"Generated PDF report: {filename} ({payload.get('engine', 'unknown')} engine, {payload.get('size', 0)} bytes)",
            path=output_path,
            url=gcs_url,
            metadata={
                **artifact_storage_metadata(session_id, run_id, output_relative_path),
                "content_type": "application/pdf",
                "size": payload.get("size", 0),
                "engine": payload.get("engine", "unknown"),
                "role": "deliverable",
            },
        )
        # Emit the durable artifact to the UI directly. The orchestrator excludes
        # self-persisting document tools from its reference-artifact path, so this
        # is the single canonical artifact for this file (one id everywhere).
        await _notify_artifact_created(artifact)

        return {
            "status": "success",
            "summary": f"Generated PDF report: {filename}",
            "detail": {
                "filename": filename,
                "path": output_path,
                "relative_path": output_relative_path,
                "artifact_id": artifact.artifact_id,
                "url": gcs_url,
                "engine": payload.get("engine"),
                "size": payload.get("size"),
            },
        }
    except Exception as e:
        logger.exception("Failed to promote PDF to artifact")
        # Durable persistence is required: a file only in the ephemeral sandbox
        # is not a deliverable. Report failure instead of faking success.
        return {
            "status": "error",
            "summary": f"PDF generated but could not be persisted to durable storage: {e}",
            "detail": {"filename": filename, "path": output_path},
            "error_code": "ARTIFACT_PERSISTENCE_FAILED",
        }


@normalized_tool
async def publish_html_artifact(
    title: str,
    html: str,
    filename: str | None = None,
) -> dict[str, Any]:
    """Publish a self-contained HTML/CSS/JS artifact directly to the UI preview panel.

    Use this for simple calculators, dashboards, interactive reports, charts, and
    single-page tools that do not require a React/Next dev server.
    """
    session_id = get_session_id()
    run_id = get_run_id()
    history_repo = get_history_repository()
    if history_repo is None:
        return {
            "status": "error",
            "summary": "Artifact history repository is not available.",
            "detail": None,
            "error_code": "ARTIFACT_HISTORY_UNAVAILABLE",
        }
    clean_title = (title or "HTML Artifact").strip()[:160]
    html_content = _ensure_full_html_document(clean_title, str(html or ""))
    if len(html_content.strip()) < 20:
        return {
            "status": "error",
            "summary": "HTML content is too short to publish.",
            "detail": None,
            "error_code": "INVALID_HTML_ARTIFACT",
        }

    output_filename = _safe_html_filename(filename or clean_title)
    relative_path = f"outputs/{output_filename}"
    content_bytes = html_content.encode("utf-8")
    gcs_url = await upload_artifact_async(
        session_id=session_id,
        run_id=run_id,
        relative_path=relative_path,
        content=html_content,
    )

    url = gcs_url
    if not url and len(content_bytes) <= _HTML_DATA_URI_LIMIT_BYTES:
        encoded = base64.b64encode(content_bytes).decode("ascii")
        url = f"data:text/html;charset=utf-8;base64,{encoded}"
    if not url:
        return {
            "status": "error",
            "summary": "Failed to store HTML artifact and it is too large for inline fallback.",
            "detail": {"filename": output_filename, "size": len(content_bytes)},
            "error_code": "HTML_ARTIFACT_STORAGE_FAILED",
        }

    metadata = {
        **artifact_storage_metadata(session_id, run_id, relative_path),
        "relative_path": relative_path,
        "content_type": "text/html; charset=utf-8",
        "size": len(content_bytes),
        "render_mode": "iframe",
        "artifact_role": "html_preview",
        "role": "deliverable",
        "storage": "gcs" if gcs_url else "data_uri",
    }
    artifact = await history_repo.create_artifact(
        session_id=session_id,
        run_id=run_id,
        kind="html",
        title=clean_title,
        preview=_html_preview_text(html_content),
        path=relative_path,
        url=url,
        metadata=metadata,
    )
    await _notify_artifact_created(artifact)

    return {
        "status": "success",
        "summary": f"Published HTML artifact: {clean_title}",
        "detail": {
            "artifact_id": artifact.artifact_id,
            "title": clean_title,
            "filename": output_filename,
            "path": relative_path,
            "url": url,
            "metadata": artifact.metadata or metadata,
        },
    }


@normalized_tool(needs_sandbox=True)
async def publish_app_preview(port: int, title: str | None = None) -> dict[str, Any]:
    """Publish a live HTTPS preview URL for an app already listening in the sandbox.

    Use after starting a Vite/Next/Flask (or similar) server bound to 0.0.0.0.
    Simple single-file HTML should use publish_html_artifact instead.

    Args:
        port: TCP port the app is listening on inside the sandbox (e.g. 5173).
        title: Optional label shown in the Preview tab.
    """
    try:
        port_number = int(port)
    except (TypeError, ValueError):
        return tool_error("port must be an integer.", error_code="INVALID_INPUT")
    if port_number < 1 or port_number > 65535:
        return tool_error("port must be between 1 and 65535.", error_code="INVALID_INPUT")

    sandbox = get_sandbox()
    if sandbox is None or not sandbox.is_alive:
        return tool_error(
            "Sandbox is not running. Start the app in this session, then retry.",
            error_code="SANDBOX_NOT_RUNNING",
            retryable=True,
        )
    if not sandbox.probe_listening_port(port_number):
        return tool_error(
            f"Nothing is listening on 127.0.0.1:{port_number}. "
            "Start the server bound to 0.0.0.0 on that port (for example "
            "`npm run dev -- --host 0.0.0.0 --port 5173`), then call this tool again.",
            error_code="PORT_NOT_LISTENING",
            suggested_alternatives=["run_command", "terminal_worker"],
        )

    try:
        url = sandbox.get_preview_url(port_number)
    except Exception as exc:
        return tool_error(
            f"Could not resolve a public preview URL for port {port_number}: {exc}",
            error_code="PREVIEW_URL_FAILED",
            retryable=True,
        )

    clean_title = (title or "App preview").strip()[:160] or "App preview"
    try:
        workspace_path = get_workspace_path()
    except RuntimeError:
        workspace_path = ""

    try:
        from nexus.tools.preview_hosts import prepare_vite_preview_for_e2b

        prep = await asyncio.to_thread(
            prepare_vite_preview_for_e2b,
            sandbox,
            port=port_number,
            workspace_path=workspace_path,
            public_url=url,
        )
        if prep.get("restarted"):
            for _ in range(12):
                await asyncio.sleep(0.5)
                if sandbox.probe_listening_port(port_number):
                    break
    except Exception:
        logger.debug("Vite preview host preparation failed", exc_info=True)

    await emit_sandbox_event(
        {
            "type": "app_preview",
            "url": url,
            "port": port_number,
            "title": clean_title,
            "workspace_path": workspace_path,
        }
    )

    return tool_success(
        f"Live preview ready: {url}",
        detail={
            "url": url,
            "port": port_number,
            "title": clean_title,
            "workspace_path": workspace_path,
            "note": "This URL stays live only while the sandbox is running.",
        },
    )


@normalized_tool(needs_sandbox=True)
async def save_as_artifact(path: str, title: str | None = None) -> dict[str, Any]:
    """
    Take any file created in the workspace and promote it to an artifact for the user to view.
    
    Args:
        path: Relative path to the file in the workspace (e.g., 'outputs/chart.png').
        title: Optional display title for the artifact.
    """
    sandbox = get_sandbox()
    session_id = get_session_id()
    run_id = get_run_id()
    history_repo = get_history_repository()

    # Reject path traversal so a promoted artifact always maps to a stable,
    # workspace-relative GCS blob (never escapes the run's namespace).
    if ".." in (path or "").replace("\\", "/").split("/"):
        return {
            "status": "error",
            "summary": f"Invalid path (traversal not allowed): {path}",
            "detail": None,
            "error_code": "INVALID_ARTIFACT_PATH",
        }

    if not sandbox.path_exists(path):
        return {
            "status": "error",
            "summary": f"File not found: {path}",
            "detail": None
        }
        
    kind = "file"
    if path.endswith((".png", ".jpg", ".jpeg")):
        kind = "image"
    elif path.endswith(".pdf"):
        kind = "pdf"
    elif path.endswith((".html", ".htm")):
        kind = "html"
    elif path.endswith((".xlsx", ".xls", ".csv")):
        kind = "spreadsheet"
    elif path.endswith((".pptx", ".ppt")):
        kind = "presentation"
    elif path.endswith(".docx"):
        kind = "document"
    elif path.endswith(".json"):
        kind = "data"
        
    try:
        content = sandbox.read_binary_file(path)
        gcs_url = await upload_artifact_async(
            session_id=session_id,
            run_id=run_id,
            relative_path=path,
            content=content
        )
        
        metadata = artifact_storage_metadata(session_id, run_id, path)
        metadata["role"] = "deliverable"
        if kind == "html":
            metadata.update({
                "content_type": "text/html; charset=utf-8",
                "render_mode": "iframe",
                "artifact_role": "html_preview",
            })
        artifact = await history_repo.create_artifact(
            session_id=session_id,
            run_id=run_id,
            kind=kind,
            title=title or os.path.basename(path),
            preview=f"Promoted {kind}: {path}",
            path=path,
            url=gcs_url,
            metadata=metadata,
        )
        await _notify_artifact_created(artifact)

        return {
            "status": "success",
            "summary": f"Promoted {path} to artifact.",
            "detail": {
                "artifact_id": artifact.artifact_id,
                "url": gcs_url,
                "metadata": artifact.metadata or {},
            }
        }
    except Exception as e:
        logger.exception("Failed to promote file to artifact")
        return {
            "status": "error",
            "summary": f"Failed to promote {path} to artifact: {str(e)}",
            "detail": None
        }

@normalized_tool(needs_sandbox=True)
async def generate_excel_report(
    title: str,
    headers: list[str],
    rows: list[list[str]],
    filename: str | None = None,
    sheet_name: str = "Report",
) -> dict[str, Any]:
    """
    Generate a styled Excel (.xlsx) workbook from tabular data.

    Use for local spreadsheets. If the user asked for a Google Sheet and Drive
    is connected, call create_drive_sheet instead.

    Args:
        title: Report title (written as the first row in bold).
        headers: Column header names.
        rows: List of row data (each row is a list of cell values).
        filename: Optional desired filename (e.g., 'sales_data.xlsx').
        sheet_name: Name of the Excel worksheet.
    """
    sandbox = get_sandbox()
    session_id = get_session_id()
    run_id = get_run_id()
    history_repo = get_history_repository()

    if not filename:
        filename = f"report_{run_id[:8]}.xlsx"
    if not filename.endswith(".xlsx"):
        filename += ".xlsx"
    filename = os.path.basename(filename)
    output_relative_path = f"outputs/{filename}"
    output_path = f"{get_workspace_path().rstrip('/')}/{output_relative_path}"

    # Write data as JSON for the sandbox script to read
    import json as _json
    data_path = f"/tmp/_xlsx_data_{run_id}.json"
    sandbox.write_text_file(data_path, _json.dumps({
        "title": title,
        "headers": headers,
        "rows": rows,
        "sheet_name": sheet_name,
    }))

    xlsx_script = _SANDBOX_DEPS_BOOTSTRAP + textwrap.dedent(f"""\
import json, os, sys
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

data = json.loads(open({repr(data_path)}).read())
out_path = {repr(output_path)}
os.makedirs(os.path.dirname(out_path), exist_ok=True)

wb = Workbook()
ws = wb.active
ws.title = data["sheet_name"]

# Title row
ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(data["headers"]), 1))
title_cell = ws.cell(row=1, column=1, value=data["title"])
title_cell.font = Font(bold=True, size=16, color="F5F5F4", name="Calibri")
title_cell.fill = PatternFill(start_color="0B0B0E", end_color="0B0B0E", fill_type="solid")
title_cell.alignment = Alignment(horizontal="left", vertical="center")
ws.row_dimensions[1].height = 28

accent = PatternFill(start_color="2DD4BF", end_color="2DD4BF", fill_type="solid")
ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=max(len(data["headers"]), 1))
ws.cell(row=2, column=1).fill = accent
ws.row_dimensions[2].height = 6

# Header row
header_fill = PatternFill(start_color="0B0B0E", end_color="0B0B0E", fill_type="solid")
header_font = Font(bold=True, color="2DD4BF", size=11, name="Calibri")
thin_border = Border(
    left=Side(style="thin", color="D6D3D1"), right=Side(style="thin", color="D6D3D1"),
    top=Side(style="thin", color="D6D3D1"), bottom=Side(style="thin", color="D6D3D1"),
)
for col_idx, header in enumerate(data["headers"], 1):
    cell = ws.cell(row=3, column=col_idx, value=header)
    cell.font = header_font
    cell.fill = header_fill
    cell.alignment = Alignment(horizontal="left", vertical="center")
    cell.border = thin_border
ws.row_dimensions[3].height = 22

def coerce(value):
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        return value
    text = str(value).strip()
    if not text:
        return ""
    compact = text.replace(",", "")
    try:
        if compact.count(".") == 1 and compact.replace(".", "", 1).lstrip("-").isdigit():
            return float(compact)
        if compact.lstrip("-").isdigit():
            return int(compact)
    except Exception:
        return value
    return value

body_font = Font(name="Calibri", size=11, color="1C1917")
alt_fill = PatternFill(start_color="FAFAF9", end_color="FAFAF9", fill_type="solid")
for row_idx, row_data in enumerate(data["rows"], 4):
    for col_idx, value in enumerate(row_data, 1):
        cell = ws.cell(row=row_idx, column=col_idx, value=coerce(value))
        cell.font = body_font
        cell.border = thin_border
        cell.alignment = Alignment(vertical="center")
        if row_idx % 2 == 0:
            cell.fill = alt_fill

ws.freeze_panes = "A4"
if data["headers"]:
    last_col = ws.cell(row=3, column=len(data["headers"])).column_letter
    ws.auto_filter.ref = f"A3:{{last_col}}{{max(3, ws.max_row)}}"
ws.sheet_view.showGridLines = False

# Auto-size columns
for col_idx in range(1, len(data["headers"]) + 1):
    max_len = max(
        (len(str(ws.cell(row=r, column=col_idx).value or "")) for r in range(3, ws.max_row + 1)),
        default=10,
    )
    ws.column_dimensions[ws.cell(row=3, column=col_idx).column_letter].width = min(max(max_len + 4, 12), 42)

wb.save(out_path)
size = os.path.getsize(out_path)
print(json.dumps({{"status": "success", "path": out_path, "size": size}}))
""")

    script_path = f"/tmp/gen_xlsx_{run_id}.py"
    sandbox.write_text_file(script_path, xlsx_script)

    res = sandbox.run_command(f"python3 {script_path}", timeout=60)
    sandbox.run_command(f"rm -f {shlex.quote(script_path)} {shlex.quote(data_path)}", timeout=10)

    if res.get("exit_code") != 0:
        return {
            "status": "error",
            "summary": f"Excel generation failed: {res.get('stderr') or res.get('stdout')}",
            "detail": res,
        }

    try:
        import json as _json2
        payload = _json2.loads(str(res.get("stdout") or "").strip())
    except Exception:
        payload = {}

    try:
        content = sandbox.read_binary_file(output_path)
        gcs_url = await upload_artifact_async(
            session_id=session_id,
            run_id=run_id,
            relative_path=output_relative_path,
            content=content,
        )

        artifact = await history_repo.create_artifact(
            session_id=session_id,
            run_id=run_id,
            kind="spreadsheet",
            title=title or filename,
            preview=f"Generated Excel report: {filename} ({len(rows)} rows)",
            path=output_path,
            url=gcs_url,
            metadata={
                **artifact_storage_metadata(session_id, run_id, output_relative_path),
                "content_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "size": payload.get("size", 0),
                "row_count": len(rows),
                "role": "deliverable",
            },
        )
        await _notify_artifact_created(artifact)

        return {
            "status": "success",
            "summary": f"Generated Excel report: {filename} ({len(rows)} rows)",
            "detail": {
                "filename": filename,
                "path": output_path,
                "relative_path": output_relative_path,
                "artifact_id": artifact.artifact_id,
                "url": gcs_url,
            },
        }
    except Exception as e:
        logger.exception("Failed to promote Excel to artifact")
        return {
            "status": "error",
            "summary": f"Excel generated but could not be persisted to durable storage: {e}",
            "detail": {"filename": filename, "path": output_path},
            "error_code": "ARTIFACT_PERSISTENCE_FAILED",
        }


@normalized_tool(needs_sandbox=True)
async def generate_docx_report(
    title: str,
    markdown_content: str,
    filename: str | None = None
) -> dict[str, Any]:
    """
    Generate a Word (.docx) document from Markdown content.

    Write the body as real markdown with H2 sections, short paragraphs, and
    lists. The tool applies professional typography; do not invent .docx bytes.

    Args:
        title: The title of the document.
        markdown_content: The body in Markdown format.
        filename: Optional desired filename (e.g., 'summary.docx').
    """
    sandbox = get_sandbox()
    session_id = get_session_id()
    run_id = get_run_id()
    history_repo = get_history_repository()

    if not filename:
        filename = f"report_{run_id[:8]}.docx"
    if not filename.endswith(".docx"):
        filename += ".docx"
    filename = os.path.basename(filename)
    output_relative_path = f"outputs/{filename}"
    output_path = f"{get_workspace_path().rstrip('/')}/{output_relative_path}"

    md_temp = f"/tmp/_docx_md_{run_id}.md"
    sandbox.write_text_file(md_temp, markdown_content)

    docx_script = _SANDBOX_DEPS_BOOTSTRAP + textwrap.dedent(f"""\
import json, os, sys, pathlib, re
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from docx.oxml.ns import qn
from docx.oxml import OxmlElement

title = {repr(title)}
md_path = {repr(md_temp)}
out_path = {repr(output_path)}

os.makedirs(os.path.dirname(out_path), exist_ok=True)
md_text = pathlib.Path(md_path).read_text(encoding="utf-8")

doc = Document()
for section in doc.sections:
    section.top_margin = Inches(1.0)
    section.bottom_margin = Inches(0.9)
    section.left_margin = Inches(1.05)
    section.right_margin = Inches(1.05)

def set_run_font(run, name="Calibri", size=11, bold=False, color=None):
    run.font.name = name
    try:
        run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    except Exception:
        pass
    run.font.size = Pt(size)
    run.bold = bold
    if color is not None:
        run.font.color.rgb = color

styles = doc.styles
normal = styles["Normal"]
normal.font.name = "Calibri"
normal.font.size = Pt(11)
normal.font.color.rgb = RGBColor(0x1C, 0x19, 0x17)

# Style the title
title_para = doc.add_heading(title, level=0)
title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
for run in title_para.runs:
    set_run_font(run, size=26, bold=True, color=RGBColor(0x0B, 0x0B, 0x0E))
accent = doc.add_paragraph()
accent_run = accent.add_run(" ")
pPr = accent._p.get_or_add_pPr()
pBdr = OxmlElement("w:pBdr")
bottom = OxmlElement("w:bottom")
bottom.set(qn("w:val"), "single")
bottom.set(qn("w:sz"), "12")
bottom.set(qn("w:space"), "1")
bottom.set(qn("w:color"), "2DD4BF")
pBdr.append(bottom)
pPr.append(pBdr)

# Parse markdown line by line
for line in md_text.splitlines():
    stripped = line.strip()
    if not stripped:
        doc.add_paragraph("")
    elif stripped.startswith("### "):
        heading = doc.add_heading(stripped[4:], level=3)
        for run in heading.runs:
            run.font.color.rgb = RGBColor(0x44, 0x40, 0x3C)
    elif stripped.startswith("## "):
        heading = doc.add_heading(stripped[3:], level=2)
        for run in heading.runs:
            run.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0E)
    elif stripped.startswith("# "):
        heading = doc.add_heading(stripped[2:], level=1)
        for run in heading.runs:
            run.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0E)
    elif stripped.startswith("- ") or stripped.startswith("* "):
        doc.add_paragraph(stripped[2:], style="List Bullet")
    elif re.match(r"^\\d+\\.\\s", stripped):
        text = re.sub(r"^\\d+\\.\\s", "", stripped)
        doc.add_paragraph(text, style="List Number")
    elif stripped.startswith("> "):
        p = doc.add_paragraph(stripped[2:])
        p.style = doc.styles["Intense Quote"] if "Intense Quote" in [s.name for s in doc.styles] else doc.styles["Quote"] if "Quote" in [s.name for s in doc.styles] else p.style
    elif stripped.startswith("```"):
        pass  # skip code fence markers
    else:
        p = doc.add_paragraph()
        # Handle inline bold/italic
        parts = re.split(r"(\\*\\*.*?\\*\\*|\\*.*?\\*)", stripped)
        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                run = p.add_run(part[2:-2])
                run.bold = True
            elif part.startswith("*") and part.endswith("*"):
                run = p.add_run(part[1:-1])
                run.italic = True
            else:
                p.add_run(part)

doc.save(out_path)
size = os.path.getsize(out_path)
print(json.dumps({{"status": "success", "path": out_path, "size": size}}))
""")

    script_path = f"/tmp/gen_docx_{run_id}.py"
    sandbox.write_text_file(script_path, docx_script)

    res = sandbox.run_command(f"python3 {script_path}", timeout=60)
    sandbox.run_command(f"rm -f {shlex.quote(script_path)}", timeout=10)

    if res.get("exit_code") != 0:
        return {
            "status": "error",
            "summary": f"DOCX generation failed: {res.get('stderr') or res.get('stdout')}",
            "detail": res,
        }

    try:
        import json as _json3
        payload = _json3.loads(str(res.get("stdout") or "").strip())
    except Exception:
        payload = {}

    try:
        content = sandbox.read_binary_file(output_path)
        gcs_url = await upload_artifact_async(
            session_id=session_id,
            run_id=run_id,
            relative_path=output_relative_path,
            content=content,
        )

        # Generate the PDF sibling for preview
        pdf_filename = filename[:-5] + ".pdf" if filename.endswith(".docx") else filename + ".pdf"
        pdf_relative_path = f"outputs/{pdf_filename}"
        pdf_output_path = f"{get_workspace_path().rstrip('/')}/{pdf_relative_path}"
        pdf_gcs_url: str | None = None
        pdf_ready = False
        try:
            pdf_script = _SANDBOX_DEPS_BOOTSTRAP + textwrap.dedent(f"""\
import json, os, sys, pathlib

title = {repr(title)}
md_path = {repr(md_temp)}
out_path = {repr(pdf_output_path)}
CSS_TEMPLATE = {json.dumps(PDF_REPORT_CSS)}

os.makedirs(os.path.dirname(out_path), exist_ok=True)
md_text = pathlib.Path(md_path).read_text(encoding="utf-8")

def try_weasyprint(md_text, title, out_path):
    import markdown2
    from weasyprint import HTML, CSS
    html_body = markdown2.markdown(
        md_text,
        extras=["fenced-code-blocks", "tables", "break-on-newline",
                "header-ids", "strike", "task_list"]
    )
    full_html = f\"\"\"<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>{{title}}</title></head>
<body>
<header class="cover"><p class="kicker">Document</p><h1>{{title}}</h1></header>
{{html_body}}
</body></html>\"\"\"
    HTML(string=full_html).write_pdf(out_path, stylesheets=[CSS(string=CSS_TEMPLATE)])
    return True

try:
    ok = try_weasyprint(md_text, title, out_path)
    print(json.dumps({{"status": "success", "size": os.path.getsize(out_path)}}))
except Exception as e:
    print(json.dumps({{"status": "error", "message": str(e)}}))
""")
            pdf_script_path = f"/tmp/gen_docx_pdf_{run_id}.py"
            sandbox.write_text_file(pdf_script_path, pdf_script)
            pdf_res = sandbox.run_command(f"python3 {pdf_script_path}", timeout=120)
            sandbox.run_command(f"rm -f {shlex.quote(pdf_script_path)}", timeout=10)
            if pdf_res.get("exit_code") == 0:
                pdf_content = sandbox.read_binary_file(pdf_output_path)
                pdf_ready = True
                pdf_gcs_url = await upload_artifact_async(
                    session_id=session_id,
                    run_id=run_id,
                    relative_path=pdf_relative_path,
                    content=pdf_content,
                )
        except Exception:
            logger.warning("Failed to generate PDF sibling for DOCX preview", exc_info=True)
        sandbox.run_command(f"rm -f {shlex.quote(md_temp)}", timeout=10)

        artifact = await history_repo.create_artifact(
            session_id=session_id,
            run_id=run_id,
            kind="document",
            title=title or filename,
            preview=f"Generated Word document: {filename}",
            path=output_path,
            url=gcs_url,
            metadata={
                **artifact_storage_metadata(session_id, run_id, output_relative_path),
                "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "size": payload.get("size", 0),
                "role": "deliverable",
                **({"preview_url": pdf_gcs_url} if pdf_gcs_url else {}),
                **({"preview_path": pdf_relative_path} if pdf_ready else {}),
                **({"preview_content_type": "application/pdf"} if pdf_ready else {}),
            },
        )
        await _notify_artifact_created(artifact)

        return {
            "status": "success",
        "summary": f"Generated Word document: {filename}",
        "detail": {
            "filename": filename,
            "path": output_path,
            "relative_path": output_relative_path,
            "artifact_id": artifact.artifact_id,
            "url": gcs_url,
            **({"preview_url": pdf_gcs_url} if pdf_gcs_url else {}),
        },
    }
    except Exception as e:
        logger.exception("Failed to promote DOCX to artifact")
        return {
            "status": "error",
            "summary": f"DOCX generated but could not be persisted to durable storage: {e}",
            "detail": {"filename": filename, "path": output_path},
            "error_code": "ARTIFACT_PERSISTENCE_FAILED",
        }


@normalized_tool(needs_sandbox=True)
async def generate_pptx_report(
    title: str,
    slides: list[dict[str, Any]] | None = None,
    filename: str | None = None,
) -> dict[str, Any]:
    """Generate a widescreen PowerPoint deck with a modern dark design system.

    The tool applies near-black backgrounds, strong type, and a single teal
    accent. You write the content and pick layouts; do not invent PPTX bytes.

    Each slide is an object:
    - layout: title | section | content | split | stats | quote | closing
    - kicker: optional short eyebrow (e.g. "Q3 Review")
    - title: headline, under 8 words
    - subtitle: optional supporting line
    - bullets: 3-5 short lines for content/closing slides
    - left / right: bullet lists for split slides
    - stats: up to 4 objects {value, label} for KPI slides
    - quote / attribution: for quote slides
    - footnote: optional footer

    Start with a title slide, use section dividers between chapters, and end
    with a closing slide. Prefer generate_pptx_report over python-pptx scripts.

    Args:
        title: Deck title used on the cover when the first slide has no kicker.
        slides: Slide objects as described above. Legacy {title, bullets} still works.
        filename: Optional desired filename (e.g. 'pitch.pptx').
    """
    sandbox = get_sandbox()
    session_id = get_session_id()
    run_id = get_run_id()
    history_repo = get_history_repository()

    normalized = _normalize_slides(slides, deck_title=title)
    if not normalized:
        return {
            "status": "error",
            "summary": "Provide at least one slide with a title or bullets.",
            "detail": None,
            "error_code": "INVALID_SLIDES",
        }
    if title.strip() and not any(item.get("layout") == "title" for item in normalized):
        normalized = [
            {
                "layout": "title",
                "kicker": "",
                "title": title.strip(),
                "subtitle": "",
                "bullets": [],
                "left": [],
                "right": [],
                "stats": [],
                "quote": "",
                "attribution": "",
                "footnote": "",
            },
            *normalized,
        ]

    if not filename:
        filename = f"deck_{run_id[:8]}.pptx"
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    filename = os.path.basename(filename)
    output_relative_path = f"outputs/{filename}"
    output_path = f"{get_workspace_path().rstrip('/')}/{output_relative_path}"
    html_filename = filename[:-5] + ".html"
    html_relative_path = f"outputs/{html_filename}"
    html_output_path = f"{get_workspace_path().rstrip('/')}/{html_relative_path}"

    import json as _json
    data_path = f"/tmp/_pptx_data_{run_id}.json"
    sandbox.write_text_file(data_path, _json.dumps({
        "title": title,
        "slides": normalized,
    }))

    generator_source = pptx_generator_source()
    bootstrap_path = f"/tmp/_pptx_deps_{run_id}.py"
    script_path = f"/tmp/gen_pptx_{run_id}.py"
    sandbox.write_text_file(bootstrap_path, _SANDBOX_DEPS_BOOTSTRAP)
    sandbox.write_text_file(script_path, generator_source)

    # #region agent log
    _dbg_compile = "ok"
    try:
        compile(generator_source, script_path, "exec")
    except SyntaxError as _dbg_exc:
        _dbg_compile = f"{_dbg_exc.msg}: line {_dbg_exc.lineno}"
    _dbg_pptx_log(
        "H1",
        "docs.py:generate_pptx_report:compose",
        "sandbox pptx generator as standalone file",
        {
            "line_count": len(generator_source.splitlines()),
            "future_import_line": next(
                (
                    i + 1
                    for i, line in enumerate(generator_source.splitlines())
                    if "from __future__ import" in line
                ),
                None,
            ),
            "compile": _dbg_compile,
            "mode": "argv",
            "slide_count": len(normalized),
            "filename": filename,
            "preview": generator_source.splitlines()[:14],
        },
    )
    # #endregion

    sandbox.run_command(f"python3 {shlex.quote(bootstrap_path)}", timeout=240)
    res = sandbox.run_command(
        "python3 "
        f"{shlex.quote(script_path)} {shlex.quote(data_path)} "
        f"{shlex.quote(output_path)} {shlex.quote(html_output_path)}",
        timeout=90,
    )
    sandbox.run_command(
        f"rm -f {shlex.quote(script_path)} {shlex.quote(data_path)} {shlex.quote(bootstrap_path)}",
        timeout=10,
    )

    # #region agent log
    _dbg_pptx_log(
        "H1",
        "docs.py:generate_pptx_report:sandbox",
        "sandbox python3 gen_pptx result",
        {
            "exit_code": res.get("exit_code"),
            "stderr": str(res.get("stderr") or "")[:800],
            "stdout": str(res.get("stdout") or "")[:400],
            "filename": filename,
            "syntax_error": "from __future__" in str(res.get("stderr") or ""),
        },
    )
    # #endregion

    if res.get("exit_code") != 0:
        stderr_text = str(res.get("stderr") or res.get("stdout") or "")
        error_code = (
            "PPTX_GENERATOR_TEMPLATE_SYNTAX_ERROR"
            if "from __future__" in stderr_text
            else "PPTX_GENERATOR_FAILED"
        )
        return {
            "status": "error",
            "summary": f"PPTX generation failed: {stderr_text}",
            "detail": res,
            "error_code": error_code,
        }

    try:
        import json as _json2
        payload = _json2.loads(str(res.get("stdout") or "").strip())
    except Exception:
        payload = {}

    try:
        content = sandbox.read_binary_file(output_path)
        gcs_url = await upload_artifact_async(
            session_id=session_id,
            run_id=run_id,
            relative_path=output_relative_path,
            content=content,
        )

        html_gcs_url: str | None = None
        html_ready = False
        try:
            html_content = sandbox.read_binary_file(html_output_path)
            html_ready = True
            html_gcs_url = await upload_artifact_async(
                session_id=session_id,
                run_id=run_id,
                relative_path=html_relative_path,
                content=html_content,
            )
        except Exception:
            logger.warning("Failed to upload HTML sibling for PPTX preview", exc_info=True)

        artifact = await history_repo.create_artifact(
            session_id=session_id,
            run_id=run_id,
            kind="presentation",
            title=title or filename,
            preview=f"Generated slide deck: {filename} ({len(normalized)} slides)",
            path=output_path,
            url=gcs_url,
            metadata={
                **artifact_storage_metadata(session_id, run_id, output_relative_path),
                "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "size": payload.get("size", 0),
                "slide_count": len(normalized),
                "role": "deliverable",
                **({"preview_url": html_gcs_url} if html_gcs_url else {}),
                **({"preview_path": html_relative_path} if html_ready else {}),
                **({"preview_content_type": "text/html; charset=utf-8"} if html_ready else {}),
                **({"render_mode": "iframe"} if html_ready else {}),
            },
        )
        await _notify_artifact_created(artifact)

        return {
            "status": "success",
            "summary": f"Generated slide deck: {filename} ({len(normalized)} slides)",
            "detail": {
                "filename": filename,
                "path": output_path,
                "relative_path": output_relative_path,
                "artifact_id": artifact.artifact_id,
                "url": gcs_url,
                **({"preview_url": html_gcs_url} if html_gcs_url else {}),
            },
        }
    except Exception as e:
        logger.exception("Failed to promote PPTX to artifact")
        return {
            "status": "error",
            "summary": f"PPTX generated but could not be persisted to durable storage: {e}",
            "detail": {"filename": filename, "path": output_path},
            "error_code": "ARTIFACT_PERSISTENCE_FAILED",
        }
