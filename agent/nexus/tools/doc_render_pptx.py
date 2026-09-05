# Copyright (c) 2026 nandan-d14. All rights reserved.
# Proprietary and non-commercial use only.

"""Sandbox-side renderer for themed PowerPoint decks.

``PPTX_RENDERER`` is executed inside the E2B sandbox as ``python3 script.py
payload.json``. It builds every slide from the blank layout with explicit
geometry, which is what makes the output look designed instead of like the
stock python-pptx template stretched to 16:9.

It also emits a matching HTML deck used for the in-app preview.
"""

from __future__ import annotations

PPTX_RENDERER = r'''
import base64
import html as html_lib
import json
import math
import mimetypes
import os
import re
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

PAYLOAD = json.loads(open(sys.argv[1], encoding="utf-8").read())
THEME = PAYLOAD["theme"]
FONTS = THEME.get("fonts") or {}
HEAD_FONT = FONTS.get("heading") or "Segoe UI"
BODY_FONT = FONTS.get("body") or "Segoe UI"
MONO_FONT = FONTS.get("mono") or "Consolas"
PALETTE = [str(c) for c in (THEME.get("chart_palette") or ["4F46E5"])]

DECK_TITLE = str(PAYLOAD.get("title") or "Presentation")
DECK_SUBTITLE = str(PAYLOAD.get("subtitle") or "")
DECK_AUTHOR = str(PAYLOAD.get("author") or "")
FOOTER_LABEL = str(PAYLOAD.get("footer") or DECK_TITLE)
OUT_PATH = PAYLOAD["out_path"]
HTML_PATH = PAYLOAD.get("html_path")
WORKSPACE = str(PAYLOAD.get("workspace") or "")

# Exact PowerPoint widescreen canvas. Inches(13.333) is 305 EMU short of this.
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
MARGIN = Inches(0.85)
CONTENT_W = SLIDE_W - 2 * MARGIN
TITLE_TOP = Inches(0.52)
RULE_TOP = Inches(1.44)
BODY_TOP = Inches(1.78)
BODY_BOTTOM = Inches(6.60)
BODY_H = BODY_BOTTOM - BODY_TOP
FOOTER_TOP = Inches(6.86)

BODY_W_IN = 13.333 - 2 * 0.85
BODY_H_IN = 6.60 - 1.78


def hexc(value, fallback="000000"):
    text = str(value or "").strip().lstrip("#").upper()
    if len(text) != 6:
        text = fallback
    try:
        return RGBColor.from_string(text)
    except Exception:
        return RGBColor.from_string(fallback)


def color(key, fallback="000000"):
    return hexc(THEME.get(key), fallback)


def mix(fg_key, bg_key, ratio):
    # PowerPoint has no simple text-opacity API, so blend toward the background
    # to get a soft tint that stays readable in both light and dark themes.
    def parts(key, fallback):
        text = str(THEME.get(key) or fallback).lstrip("#").upper()
        if len(text) != 6:
            text = fallback
        return [int(text[i:i + 2], 16) for i in (0, 2, 4)]

    fg = parts(fg_key, "4F46E5")
    bg = parts(bg_key, "FFFFFF")
    blended = "".join(
        "{:02X}".format(int(round(fg[i] * ratio + bg[i] * (1 - ratio)))) for i in range(3)
    )
    return hexc(blended)


# ---------------------------------------------------------------- primitives

def kill_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def add_rect(slide, left, top, width, height, fill=None, alpha=None,
             shape_type=MSO_SHAPE.RECTANGLE, radius=None, line_color=None,
             line_width=Pt(1)):
    shape = slide.shapes.add_shape(shape_type, int(left), int(top), int(width), int(height))
    kill_shadow(shape)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hexc(fill)
        if alpha is not None:
            set_alpha(shape, alpha)
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hexc(line_color)
        shape.line.width = line_width
    if radius is not None:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    try:
        shape.text_frame.word_wrap = True
    except Exception:
        pass
    return shape


def set_alpha(shape, percent):
    # python-pptx has no transparency API; patch the solidFill element directly.
    try:
        from pptx.oxml.ns import qn

        spPr = shape.fill._xPr
        solid = spPr.find(qn("a:solidFill"))
        if solid is None:
            return
        srgb = solid.find(qn("a:srgbClr"))
        if srgb is None:
            return
        node = srgb.makeelement(qn("a:alpha"), {"val": str(int(max(0, min(100, percent)) * 1000))})
        srgb.append(node)
    except Exception:
        pass


def add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


INLINE_RE = re.compile(r"(\*\*.+?\*\*|__.+?__|\*[^*]+?\*|`[^`]+?`)")


def add_runs(paragraph, text, size, ink, font, bold=False, italic=False, mono_color=None):
    pieces = [p for p in INLINE_RE.split(str(text)) if p]
    if not pieces:
        pieces = [""]
    for piece in pieces:
        run = paragraph.add_run()
        is_bold, is_italic, is_mono = bold, italic, False
        if (piece.startswith("**") and piece.endswith("**") and len(piece) > 4) or (
            piece.startswith("__") and piece.endswith("__") and len(piece) > 4
        ):
            run.text = piece[2:-2]
            is_bold = True
        elif piece.startswith("*") and piece.endswith("*") and len(piece) > 2:
            run.text = piece[1:-1]
            is_italic = True
        elif piece.startswith("`") and piece.endswith("`") and len(piece) > 2:
            run.text = piece[1:-1]
            is_mono = True
        else:
            run.text = piece
        fnt = run.font
        fnt.size = Pt(size)
        fnt.bold = is_bold
        fnt.italic = is_italic
        fnt.name = MONO_FONT if is_mono else font
        fnt.color.rgb = hexc(mono_color) if (is_mono and mono_color) else ink


def write_line(tf, text, size, ink, font=None, bold=False, italic=False,
               align=PP_ALIGN.LEFT, space_before=0, space_after=0,
               line_spacing=None, first=False, level=0, bullet=None,
               mono_color=None, rich=True):
    paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
    if bullet is not None:
        apply_bullet(paragraph, bullet[0], bullet[1], indent_in=bullet[2])
    else:
        apply_no_bullet(paragraph)
    paragraph.alignment = align
    paragraph.level = level
    if space_before:
        paragraph.space_before = Pt(space_before)
    if space_after:
        paragraph.space_after = Pt(space_after)
    if line_spacing:
        paragraph.line_spacing = line_spacing
    if rich:
        add_runs(paragraph, text, size, ink, font or BODY_FONT, bold=bold,
                 italic=italic, mono_color=mono_color)
    else:
        run = paragraph.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font or BODY_FONT
        run.font.color.rgb = ink
    return paragraph


def _clear_bullet_nodes(pPr, qn):
    for tag in ("a:buClrTx", "a:buClr", "a:buSzTx", "a:buSzPct", "a:buSzPts",
                "a:buFontTx", "a:buFont", "a:buNone", "a:buAutoNum", "a:buChar"):
        for node in pPr.findall(qn(tag)):
            pPr.remove(node)


def apply_bullet(paragraph, char, rgb, indent_in=0.30):
    # Bullet glyph + color live in pPr. Append order must follow the OOXML
    # sequence buClr -> buSzPct -> buFont -> buChar, and this runs before any
    # spacing/defRPr is set so python-pptx inserts those in valid positions.
    try:
        from pptx.oxml.ns import qn

        pPr = paragraph._p.get_or_add_pPr()
        _clear_bullet_nodes(pPr, qn)
        pPr.set("marL", str(int(Inches(indent_in))))
        pPr.set("indent", str(-int(Inches(indent_in))))
        nodes = [
            ("a:buClr", None),
            ("a:buSzPct", {"val": "88000"}),
            ("a:buFont", {"typeface": "Arial"}),
            ("a:buChar", {"char": char}),
        ]
        for tag, attrs in nodes:
            node = pPr.makeelement(qn(tag), attrs or {})
            if tag == "a:buClr":
                node.append(pPr.makeelement(qn("a:srgbClr"), {"val": str(rgb)}))
            pPr.append(node)
    except Exception:
        pass


def apply_no_bullet(paragraph):
    try:
        from pptx.oxml.ns import qn

        pPr = paragraph._p.get_or_add_pPr()
        _clear_bullet_nodes(pPr, qn)
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))
    except Exception:
        pass


# ------------------------------------------------------------------ fitting

def wrapped_lines(text, size_pt, width_in):
    # Rough advance width for a proportional face is ~0.5em per glyph.
    per_line = max(int((width_in * 72.0) / (size_pt * 0.50)), 8)
    return max(1, math.ceil(len(str(text)) / per_line))


def autosize(entries, width_in, height_in, start=22, minimum=11, gap_ratio=0.55):
    # entries: list of (text, level). Shrink until the block fits the box.
    size = start
    while size > minimum:
        total = 0.0
        for text, level in entries:
            indent = 0.35 * level
            lines = wrapped_lines(text, size, max(width_in - indent, 2.0))
            total += lines * size * 1.28 + size * gap_ratio
        if total <= height_in * 72.0:
            return size
        size -= 1
    return minimum


# ------------------------------------------------------------------- chrome

def paint_background(slide, fill):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=fill)


def slide_title(slide, text, eyebrow=""):
    top = TITLE_TOP
    if eyebrow:
        tf = add_textbox(slide, MARGIN, Inches(0.42), CONTENT_W, Inches(0.3))
        write_line(tf, eyebrow.upper(), 12, color("accent"), font=BODY_FONT,
                   bold=True, first=True, rich=False)
        top = Inches(0.82)
    size = 30 if len(str(text)) <= 60 else (26 if len(str(text)) <= 95 else 22)
    tf = add_textbox(slide, MARGIN, top, CONTENT_W, Inches(1.0))
    write_line(tf, text, size, color("ink"), font=HEAD_FONT, bold=True, first=True)
    add_rect(slide, MARGIN, RULE_TOP, Inches(1.9), Inches(0.055), fill=THEME.get("accent"))


def slide_footer(slide, index, total):
    add_rect(slide, MARGIN, Inches(6.74), CONTENT_W, Emu(9525), fill=THEME.get("hairline"))
    tf = add_textbox(slide, MARGIN, FOOTER_TOP, CONTENT_W * 0.7, Inches(0.3))
    write_line(tf, FOOTER_LABEL, 10, color("muted"), font=BODY_FONT, first=True, rich=False)
    tf2 = add_textbox(slide, MARGIN + CONTENT_W * 0.7, FOOTER_TOP, CONTENT_W * 0.3, Inches(0.3))
    write_line(tf2, str(index) + " / " + str(total), 10, color("muted"),
               font=BODY_FONT, align=PP_ALIGN.RIGHT, first=True, rich=False)


def attach_notes(slide, notes):
    text = str(notes or "").strip()
    if not text:
        return
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


# ------------------------------------------------------------------ layouts

def normalize_bullets(raw):
    out = []
    for item in raw or []:
        if isinstance(item, dict):
            text = str(item.get("text") or item.get("title") or "").strip()
            level = int(item.get("level") or 0)
        else:
            raw_text = str(item)
            stripped = raw_text.lstrip()
            indent = len(raw_text) - len(stripped)
            level = 1 if indent >= 2 else 0
            text = stripped
            while text[:2] in ("- ", "* "):
                text = text[2:].lstrip()
        if text:
            out.append((text, max(0, min(level, 2))))
    return out


BULLET_GLYPHS = ["\u25aa", "\u2013", "\u00b7"]


def render_bullets(slide, entries, left, top, width, height):
    if not entries:
        return
    width_in = width / 914400.0
    height_in = height / 914400.0
    size = autosize(entries, width_in, height_in)
    tf = add_textbox(slide, left, top, width, height)
    accent = str(THEME.get("accent") or "4F46E5").lstrip("#").upper()
    muted = str(THEME.get("muted") or "64748B").lstrip("#").upper()
    for idx, (text, level) in enumerate(entries):
        glyph = BULLET_GLYPHS[min(level, 2)]
        rgb = accent if level == 0 else muted
        ink = color("ink") if level == 0 else color("muted")
        write_line(
            tf, text,
            size if level == 0 else max(size - 2, 10),
            ink,
            font=BODY_FONT,
            first=(idx == 0),
            level=level,
            space_after=size * 0.5,
            line_spacing=1.18,
            bullet=(glyph, rgb, 0.30 + 0.25 * level),
            mono_color=THEME.get("accent"),
        )


def layout_cover(slide, item, index, total):
    paint_background(slide, THEME.get("cover_background"))
    add_rect(slide, SLIDE_W - Inches(4.6), -Inches(1.4), Inches(5.6), Inches(5.6),
             fill=THEME.get("cover_accent"), alpha=12, shape_type=MSO_SHAPE.OVAL)
    add_rect(slide, SLIDE_W - Inches(2.3), Inches(3.4), Inches(3.6), Inches(3.6),
             fill=THEME.get("cover_accent"), alpha=8, shape_type=MSO_SHAPE.OVAL)

    eyebrow = str(item.get("eyebrow") or "").strip()
    top = Inches(2.35)
    if eyebrow:
        tf = add_textbox(slide, MARGIN, Inches(1.95), Inches(9.5), Inches(0.35))
        write_line(tf, eyebrow.upper(), 13, color("cover_accent"), font=BODY_FONT,
                   bold=True, first=True, rich=False)

    add_rect(slide, MARGIN, top, Inches(2.2), Inches(0.09), fill=THEME.get("cover_accent"))

    title = str(item.get("title") or DECK_TITLE)
    size = 46 if len(title) <= 46 else (38 if len(title) <= 80 else 31)
    tf = add_textbox(slide, MARGIN, top + Inches(0.35), Inches(10.2), Inches(2.1))
    write_line(tf, title, size, color("cover_ink"), font=HEAD_FONT, bold=True,
               first=True, line_spacing=1.06)

    subtitle = str(item.get("subtitle") or DECK_SUBTITLE or "").strip()
    if subtitle:
        tf = add_textbox(slide, MARGIN, Inches(4.75), Inches(9.4), Inches(1.0))
        write_line(tf, subtitle, 19, color("cover_muted"), font=BODY_FONT,
                   first=True, line_spacing=1.25)

    meta = [str(v).strip() for v in (item.get("meta") or []) if str(v).strip()]
    if not meta:
        meta = [v for v in (DECK_AUTHOR, str(item.get("date") or PAYLOAD.get("date") or "")) if v]
    if meta:
        tf = add_textbox(slide, MARGIN, Inches(6.35), Inches(10.0), Inches(0.4))
        write_line(tf, "   \u00b7   ".join(meta), 12, color("cover_muted"),
                   font=BODY_FONT, first=True, rich=False)


def layout_section(slide, item, index, total):
    paint_background(slide, THEME.get("cover_background"))
    add_rect(slide, 0, 0, Inches(0.22), SLIDE_H, fill=THEME.get("cover_accent"))
    number = str(item.get("number") or "").strip()
    if number:
        tf = add_textbox(slide, MARGIN, Inches(1.75), Inches(4.0), Inches(1.6))
        write_line(tf, number, 68, color("cover_accent"), font=HEAD_FONT, bold=True,
                   first=True, rich=False)
    tf = add_textbox(slide, MARGIN, Inches(3.05), Inches(10.6), Inches(1.6),
                     anchor=MSO_ANCHOR.MIDDLE)
    title = str(item.get("title") or "Section")
    write_line(tf, title, 40 if len(title) <= 52 else 32, color("cover_ink"),
               font=HEAD_FONT, bold=True, first=True, line_spacing=1.08)
    body = str(item.get("body") or item.get("subtitle") or "").strip()
    if body:
        tf = add_textbox(slide, MARGIN, Inches(4.65), Inches(9.4), Inches(1.2))
        write_line(tf, body, 17, color("cover_muted"), font=BODY_FONT, first=True,
                   line_spacing=1.25)
    slide_footer(slide, index, total)


def layout_bullets(slide, item, index, total):
    paint_background(slide, THEME.get("background"))
    slide_title(slide, item.get("title") or "", item.get("eyebrow") or "")
    entries = normalize_bullets(item.get("bullets"))
    body_text = str(item.get("body") or "").strip()
    top = BODY_TOP
    height = BODY_H
    if body_text:
        tf = add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.9))
        write_line(tf, body_text, 16, color("muted"), font=BODY_FONT, first=True,
                   line_spacing=1.3)
        top = BODY_TOP + Inches(0.95)
        height = BODY_BOTTOM - top
    render_bullets(slide, entries, MARGIN, top, CONTENT_W, height)
    slide_footer(slide, index, total)


def layout_columns(slide, item, index, total):
    paint_background(slide, THEME.get("background"))
    slide_title(slide, item.get("title") or "", item.get("eyebrow") or "")
    columns = [c for c in (item.get("columns") or []) if isinstance(c, dict)][:3]
    if not columns:
        layout_bullets(slide, item, index, total)
        return
    gap = Inches(0.36)
    count = len(columns)
    card_w = int((CONTENT_W - gap * (count - 1)) / count)
    for i, column in enumerate(columns):
        left = MARGIN + i * (card_w + gap)
        add_rect(slide, left, BODY_TOP, card_w, BODY_H, fill=THEME.get("surface"),
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.045)
        add_rect(slide, left, BODY_TOP, card_w, Inches(0.075), fill=PALETTE[i % len(PALETTE)])
        pad = Inches(0.32)
        heading = str(column.get("heading") or column.get("title") or "").strip()
        inner_top = BODY_TOP + Inches(0.42)
        if heading:
            tf = add_textbox(slide, left + pad, inner_top, card_w - 2 * pad, Inches(0.75))
            write_line(tf, heading, 18, color("ink"), font=HEAD_FONT, bold=True,
                       first=True, line_spacing=1.12)
            inner_top = inner_top + Inches(0.82)
        entries = normalize_bullets(column.get("bullets") or column.get("points"))
        body = str(column.get("body") or "").strip()
        if body and not entries:
            tf = add_textbox(slide, left + pad, inner_top, card_w - 2 * pad,
                             BODY_TOP + BODY_H - inner_top - Inches(0.25))
            write_line(tf, body, 15, color("muted"), font=BODY_FONT, first=True,
                       line_spacing=1.32)
        else:
            render_bullets(slide, entries, left + pad, inner_top, card_w - 2 * pad,
                           BODY_TOP + BODY_H - inner_top - Inches(0.25))
    slide_footer(slide, index, total)


def layout_stats(slide, item, index, total):
    paint_background(slide, THEME.get("background"))
    slide_title(slide, item.get("title") or "", item.get("eyebrow") or "")
    stats = [s for s in (item.get("stats") or []) if isinstance(s, dict)][:4]
    if not stats:
        layout_bullets(slide, item, index, total)
        return
    gap = Inches(0.34)
    count = len(stats)
    card_w = int((CONTENT_W - gap * (count - 1)) / count)
    card_h = Inches(2.5)
    card_top = BODY_TOP + Inches(0.25)
    for i, stat in enumerate(stats):
        left = MARGIN + i * (card_w + gap)
        add_rect(slide, left, card_top, card_w, card_h, fill=THEME.get("surface"),
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        add_rect(slide, left, card_top, Inches(0.075), card_h, fill=PALETTE[i % len(PALETTE)])
        pad = Inches(0.34)
        value = str(stat.get("value") or stat.get("number") or "").strip()
        vsize = 42 if len(value) <= 5 else (34 if len(value) <= 8 else 26)
        tf = add_textbox(slide, left + pad, card_top + Inches(0.36), card_w - 2 * pad, Inches(0.9))
        write_line(tf, value, vsize, hexc(PALETTE[i % len(PALETTE)]), font=HEAD_FONT,
                   bold=True, first=True, rich=False)
        label = str(stat.get("label") or "").strip()
        if label:
            tf = add_textbox(slide, left + pad, card_top + Inches(1.28), card_w - 2 * pad, Inches(0.55))
            write_line(tf, label, 14, color("ink"), font=BODY_FONT, bold=True,
                       first=True, line_spacing=1.15)
        caption = str(stat.get("caption") or stat.get("detail") or "").strip()
        if caption:
            tf = add_textbox(slide, left + pad, card_top + Inches(1.82), card_w - 2 * pad, Inches(0.62))
            write_line(tf, caption, 11, color("muted"), font=BODY_FONT, first=True,
                       line_spacing=1.2)
    entries = normalize_bullets(item.get("bullets"))
    if entries:
        render_bullets(slide, entries[:3], MARGIN, card_top + card_h + Inches(0.34),
                       CONTENT_W, BODY_BOTTOM - (card_top + card_h + Inches(0.34)))
    slide_footer(slide, index, total)


def layout_quote(slide, item, index, total):
    paint_background(slide, THEME.get("surface"))
    add_rect(slide, 0, 0, Inches(0.2), SLIDE_H, fill=THEME.get("accent"))
    tf = add_textbox(slide, MARGIN + Inches(0.3), Inches(1.05), Inches(3.0), Inches(2.0))
    write_line(tf, "\u201c", 96, mix("accent", "surface", 0.26), font=HEAD_FONT, bold=True,
               first=True, rich=False)
    quote = str(item.get("quote") or item.get("body") or item.get("title") or "").strip()
    size = 30 if len(quote) <= 140 else (25 if len(quote) <= 240 else 20)
    tf = add_textbox(slide, MARGIN + Inches(0.35), Inches(2.25), Inches(10.6), Inches(2.9),
                     anchor=MSO_ANCHOR.MIDDLE)
    write_line(tf, quote, size, color("ink"), font=HEAD_FONT, italic=True, first=True,
               line_spacing=1.28)
    attribution = str(item.get("attribution") or item.get("author") or "").strip()
    if attribution:
        add_rect(slide, MARGIN + Inches(0.35), Inches(5.42), Inches(0.9), Inches(0.045),
                 fill=THEME.get("accent"))
        tf = add_textbox(slide, MARGIN + Inches(0.35), Inches(5.62), Inches(9.0), Inches(0.5))
        write_line(tf, attribution, 14, color("muted"), font=BODY_FONT, first=True, rich=False)
    slide_footer(slide, index, total)


CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "hbar": XL_CHART_TYPE.BAR_CLUSTERED,
    "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}


def to_number(value):
    if isinstance(value, (int, float)):
        return float(value)
    text = str(value or "").strip().replace(",", "").replace("$", "").replace("%", "")
    try:
        return float(text)
    except Exception:
        return 0.0


def normalize_chart(raw):
    if not isinstance(raw, dict):
        return None
    categories = [str(c) for c in (raw.get("categories") or raw.get("labels") or [])]
    series_raw = raw.get("series") or []
    series = []
    if isinstance(series_raw, dict):
        series_raw = [{"name": k, "values": v} for k, v in series_raw.items()]
    for entry in series_raw:
        if isinstance(entry, dict):
            name = str(entry.get("name") or entry.get("label") or "Series")
            values = [to_number(v) for v in (entry.get("values") or entry.get("data") or [])]
        else:
            continue
        if values:
            series.append({"name": name, "values": values})
    if not series and raw.get("values"):
        series = [{"name": str(raw.get("name") or "Value"),
                   "values": [to_number(v) for v in raw.get("values")]}]
    if not categories or not series:
        return None
    width = max(len(s["values"]) for s in series)
    if len(categories) < width:
        categories = categories + [""] * (width - len(categories))
    categories = categories[:width]
    for s in series:
        s["values"] = (s["values"] + [0.0] * width)[:width]
    kind = str(raw.get("type") or raw.get("kind") or "bar").strip().lower()
    return {
        "type": kind if kind in CHART_TYPES else "bar",
        "categories": categories,
        "series": series,
        "title": str(raw.get("title") or "").strip(),
    }


def draw_chart(slide, spec, left, top, width, height):
    chart_data = CategoryChartData()
    chart_data.categories = spec["categories"]
    for entry in spec["series"]:
        chart_data.add_series(entry["name"], tuple(entry["values"]))
    frame = slide.shapes.add_chart(CHART_TYPES[spec["type"]], int(left), int(top),
                                   int(width), int(height), chart_data)
    chart = frame.chart
    try:
        chart.font.size = Pt(12)
        chart.font.name = BODY_FONT
        chart.font.color.rgb = color("muted")
    except Exception:
        pass
    multi = len(spec["series"]) > 1
    is_circular = spec["type"] in ("pie", "doughnut")
    try:
        chart.has_legend = multi or is_circular
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
    except Exception:
        pass
    try:
        if is_circular:
            points = chart.plots[0].points
            for i, point in enumerate(points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = hexc(PALETTE[i % len(PALETTE)])
                point.format.line.color.rgb = color("background", "FFFFFF")
                point.format.line.width = Pt(1.5)
        else:
            for i, series in enumerate(chart.series):
                shade = hexc(PALETTE[i % len(PALETTE)])
                if spec["type"] in ("line",):
                    series.format.line.color.rgb = shade
                    series.format.line.width = Pt(2.5)
                    series.smooth = False
                else:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = shade
                    series.format.line.fill.background()
    except Exception:
        pass
    try:
        if not is_circular:
            value_axis = chart.value_axis
            value_axis.has_major_gridlines = True
            value_axis.major_gridlines.format.line.color.rgb = color("hairline")
            value_axis.major_gridlines.format.line.width = Pt(0.75)
            value_axis.format.line.fill.background()
            category_axis = chart.category_axis
            category_axis.has_major_gridlines = False
            category_axis.format.line.color.rgb = color("hairline")
    except Exception:
        pass
    return chart


def layout_chart(slide, item, index, total):
    paint_background(slide, THEME.get("background"))
    slide_title(slide, item.get("title") or "", item.get("eyebrow") or "")
    spec = normalize_chart(item.get("chart"))
    if spec is None:
        layout_bullets(slide, item, index, total)
        return
    entries = normalize_bullets(item.get("bullets"))
    if entries:
        chart_w = int(CONTENT_W * 0.62)
        draw_chart(slide, spec, MARGIN, BODY_TOP, chart_w, BODY_H)
        side_left = MARGIN + chart_w + Inches(0.4)
        side_w = int(SLIDE_W - MARGIN - side_left)
        add_rect(slide, side_left, BODY_TOP, side_w, BODY_H, fill=THEME.get("surface"),
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        render_bullets(slide, entries, side_left + Inches(0.3), BODY_TOP + Inches(0.34),
                       side_w - Inches(0.6), BODY_H - Inches(0.68))
    else:
        draw_chart(slide, spec, MARGIN, BODY_TOP, CONTENT_W, BODY_H)
    slide_footer(slide, index, total)


def normalize_table(raw):
    if not isinstance(raw, dict):
        return None
    headers = [str(h) for h in (raw.get("headers") or raw.get("columns") or [])]
    rows = []
    for row in raw.get("rows") or []:
        if isinstance(row, (list, tuple)):
            rows.append([str(c) if c is not None else "" for c in row])
        elif isinstance(row, dict) and headers:
            rows.append([str(row.get(h, "")) for h in headers])
    if not headers and rows:
        headers = ["Column " + str(i + 1) for i in range(len(rows[0]))]
    if not headers:
        return None
    width = len(headers)
    rows = [(row + [""] * width)[:width] for row in rows][:12]
    return {"headers": headers, "rows": rows}


def layout_table(slide, item, index, total):
    paint_background(slide, THEME.get("background"))
    slide_title(slide, item.get("title") or "", item.get("eyebrow") or "")
    spec = normalize_table(item.get("table"))
    if spec is None:
        layout_bullets(slide, item, index, total)
        return
    headers = spec["headers"]
    rows = spec["rows"]
    row_count = len(rows) + 1
    header_h = Inches(0.5)
    available = BODY_H - header_h
    row_h = int(min(Inches(0.46), available / max(len(rows), 1))) if rows else int(header_h)
    table_h = int(header_h + row_h * len(rows))
    shape = slide.shapes.add_table(row_count, len(headers), int(MARGIN), int(BODY_TOP),
                                   int(CONTENT_W), table_h)
    table = shape.table
    try:
        table.first_row = True
        table.horz_banding = False
    except Exception:
        pass
    table.rows[0].height = int(header_h)
    for r in range(1, row_count):
        table.rows[r].height = int(row_h)

    weights = []
    for c in range(len(headers)):
        longest = len(headers[c])
        for row in rows:
            longest = max(longest, len(row[c]))
        weights.append(max(longest, 6))
    weight_total = float(sum(weights))
    for c, weight in enumerate(weights):
        table.columns[c].width = int(CONTENT_W * (weight / weight_total))

    body_size = 12 if len(rows) <= 8 else 10
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = color("accent")
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.14)
        cell.margin_right = Inches(0.14)
        tf = cell.text_frame
        tf.word_wrap = True
        write_line(tf, header, body_size + 1, hexc("FFFFFF"), font=BODY_FONT,
                   bold=True, first=True, rich=False)
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = color("surface") if r % 2 == 0 else color("background", "FFFFFF")
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.14)
            cell.margin_right = Inches(0.14)
            tf = cell.text_frame
            tf.word_wrap = True
            write_line(tf, value, body_size, color("ink"), font=BODY_FONT, first=True, rich=False)
    slide_footer(slide, index, total)


def resolve_image(path):
    candidate = str(path or "").strip()
    if not candidate:
        return None
    options = [candidate]
    if not candidate.startswith("/") and WORKSPACE:
        options.append(os.path.join(WORKSPACE, candidate))
    for option in options:
        if os.path.isfile(option):
            return option
    return None


def image_aspect(path, default=1.6):
    try:
        from PIL import Image

        with Image.open(path) as img:
            width, height = img.size
        if height:
            return float(width) / float(height)
    except Exception:
        pass
    return default


def layout_image(slide, item, index, total):
    paint_background(slide, THEME.get("background"))
    slide_title(slide, item.get("title") or "", item.get("eyebrow") or "")
    path = resolve_image(item.get("image") or item.get("image_path"))
    caption = str(item.get("caption") or "").strip()
    entries = normalize_bullets(item.get("bullets"))
    if path is None:
        layout_bullets(slide, item, index, total)
        return

    if entries:
        box_l, box_w = MARGIN, int(CONTENT_W * 0.55)
        render_bullets(slide, entries, MARGIN + box_w + Inches(0.45), BODY_TOP + Inches(0.1),
                       int(CONTENT_W - box_w - Inches(0.45)), BODY_H - Inches(0.2))
    else:
        box_l, box_w = MARGIN, int(CONTENT_W)
    box_t = BODY_TOP
    box_h = int(BODY_H - (Inches(0.45) if caption else 0))

    aspect = image_aspect(path)
    width = box_w
    height = int(width / aspect)
    if height > box_h:
        height = box_h
        width = int(height * aspect)
    left = int(box_l + (box_w - width) / 2)
    top = int(box_t + (box_h - height) / 2)
    add_rect(slide, left - Inches(0.06), top - Inches(0.06), width + Inches(0.12),
             height + Inches(0.12), fill=THEME.get("surface"),
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.02)
    try:
        slide.shapes.add_picture(path, left, top, width, height)
    except Exception:
        pass
    if caption:
        tf = add_textbox(slide, box_l, BODY_BOTTOM - Inches(0.35), box_w, Inches(0.4))
        write_line(tf, caption, 12, color("muted"), font=BODY_FONT,
                   align=PP_ALIGN.CENTER, first=True)
    slide_footer(slide, index, total)


def layout_closing(slide, item, index, total):
    paint_background(slide, THEME.get("cover_background"))
    add_rect(slide, -Inches(1.6), SLIDE_H - Inches(3.6), Inches(5.4), Inches(5.4),
             fill=THEME.get("cover_accent"), alpha=10, shape_type=MSO_SHAPE.OVAL)
    title = str(item.get("title") or "Thank you")
    tf = add_textbox(slide, MARGIN, Inches(2.7), Inches(10.5), Inches(1.4),
                     anchor=MSO_ANCHOR.MIDDLE)
    write_line(tf, title, 44 if len(title) <= 40 else 34, color("cover_ink"),
               font=HEAD_FONT, bold=True, first=True)
    add_rect(slide, MARGIN, Inches(4.25), Inches(2.2), Inches(0.09), fill=THEME.get("cover_accent"))
    entries = normalize_bullets(item.get("bullets"))
    lines = [text for text, _ in entries]
    body = str(item.get("body") or item.get("subtitle") or "").strip()
    if body:
        lines.insert(0, body)
    if lines:
        tf = add_textbox(slide, MARGIN, Inches(4.6), Inches(9.6), Inches(1.6))
        for i, line in enumerate(lines[:4]):
            write_line(tf, line, 16, color("cover_muted"), font=BODY_FONT,
                       first=(i == 0), space_after=8, line_spacing=1.25)


LAYOUTS = {
    "cover": layout_cover,
    "title": layout_cover,
    "section": layout_section,
    "divider": layout_section,
    "bullets": layout_bullets,
    "content": layout_bullets,
    "columns": layout_columns,
    "two_column": layout_columns,
    "comparison": layout_columns,
    "stats": layout_stats,
    "kpi": layout_stats,
    "metrics": layout_stats,
    "quote": layout_quote,
    "chart": layout_chart,
    "table": layout_table,
    "image": layout_image,
    "closing": layout_closing,
    "end": layout_closing,
    "thanks": layout_closing,
}


def infer_layout(item, position, count):
    explicit = str(item.get("layout") or item.get("type") or "").strip().lower().replace("-", "_")
    if explicit in LAYOUTS:
        return explicit
    if item.get("stats"):
        return "stats"
    if item.get("chart"):
        return "chart"
    if item.get("table"):
        return "table"
    if item.get("image") or item.get("image_path"):
        return "image"
    if item.get("quote"):
        return "quote"
    if item.get("columns"):
        return "columns"
    if position == 0 and not item.get("bullets"):
        return "cover"
    return "bullets"


# --------------------------------------------------------------------- build

slides_in = [s for s in (PAYLOAD.get("slides") or []) if isinstance(s, dict)]
if not slides_in:
    slides_in = [{"layout": "cover", "title": DECK_TITLE, "subtitle": DECK_SUBTITLE}]

prepared = []
for position, item in enumerate(slides_in):
    prepared.append((infer_layout(item, position, len(slides_in)), item))

if PAYLOAD.get("auto_cover") and prepared[0][0] != "cover":
    prepared.insert(0, ("cover", {"title": DECK_TITLE, "subtitle": DECK_SUBTITLE}))

section_counter = 0
for i, (name, item) in enumerate(prepared):
    if name == "section" and not str(item.get("number") or "").strip():
        section_counter += 1
        item["number"] = "{:02d}".format(section_counter)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
try:
    blank = prs.slide_layouts[6]
except Exception:
    blank = prs.slide_layouts[len(prs.slide_layouts) - 1]

total = len(prepared)
for i, (name, item) in enumerate(prepared):
    slide = prs.slides.add_slide(blank)
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            shape._element.getparent().remove(shape._element)
    try:
        LAYOUTS[name](slide, item, i + 1, total)
    except Exception as exc:
        try:
            layout_bullets(slide, item, i + 1, total)
        except Exception:
            pass
    attach_notes(slide, item.get("notes") or item.get("speaker_notes"))

try:
    prs.core_properties.title = DECK_TITLE
    if DECK_AUTHOR:
        prs.core_properties.author = DECK_AUTHOR
    prs.core_properties.comments = "Generated by CoComputer"
except Exception:
    pass

os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
prs.save(OUT_PATH)


# ---------------------------------------------------------------- html deck

def esc(value):
    return html_lib.escape(str(value or ""))


def rich_html(text):
    out = esc(text)
    out = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", out)
    out = re.sub(r"`([^`]+?)`", r"<code>\1</code>", out)
    out = re.sub(r"(?<![\*\w])\*([^*]+?)\*(?!\w)", r"<em>\1</em>", out)
    return out


def bullets_html(entries, tag="ul"):
    if not entries:
        return ""
    items = "".join(
        '<li class="lvl' + str(level) + '">' + rich_html(text) + "</li>"
        for text, level in entries
    )
    return "<" + tag + ' class="bullets">' + items + "</" + tag + ">"


def data_uri(path):
    try:
        if os.path.getsize(path) > 3_000_000:
            return None
        mime = mimetypes.guess_type(path)[0] or "image/png"
        with open(path, "rb") as handle:
            return "data:" + mime + ";base64," + base64.b64encode(handle.read()).decode("ascii")
    except Exception:
        return None


def svg_chart(spec):
    width, height = 720, 380
    pad_l, pad_b, pad_t, pad_r = 54, 42, 18, 14
    plot_w = width - pad_l - pad_r
    plot_h = height - pad_t - pad_b
    categories = spec["categories"]
    series = spec["series"]
    parts = ['<svg viewBox="0 0 ' + str(width) + " " + str(height) + '" class="chart" preserveAspectRatio="xMidYMid meet">']

    if spec["type"] in ("pie", "doughnut"):
        values = series[0]["values"]
        total_value = sum(abs(v) for v in values) or 1.0
        cx, cy, radius = width / 2.0, height / 2.0 - 6, min(plot_h, plot_w) / 2.4
        angle = -math.pi / 2
        for i, value in enumerate(values):
            sweep = (abs(value) / total_value) * 2 * math.pi
            x1, y1 = cx + radius * math.cos(angle), cy + radius * math.sin(angle)
            angle += sweep
            x2, y2 = cx + radius * math.cos(angle), cy + radius * math.sin(angle)
            large = 1 if sweep > math.pi else 0
            d = ("M " + str(round(cx, 2)) + " " + str(round(cy, 2)) +
                 " L " + str(round(x1, 2)) + " " + str(round(y1, 2)) +
                 " A " + str(round(radius, 2)) + " " + str(round(radius, 2)) + " 0 " +
                 str(large) + " 1 " + str(round(x2, 2)) + " " + str(round(y2, 2)) + " Z")
            parts.append('<path d="' + d + '" fill="#' + PALETTE[i % len(PALETTE)] + '" stroke="var(--bg)" stroke-width="2"/>')
        if spec["type"] == "doughnut":
            parts.append('<circle cx="' + str(round(cx, 2)) + '" cy="' + str(round(cy, 2)) +
                         '" r="' + str(round(radius * 0.55, 2)) + '" fill="var(--bg)"/>')
        parts.append("</svg>")
        legend = "".join(
            '<span class="key"><i style="background:#' + PALETTE[i % len(PALETTE)] + '"></i>' +
            esc(categories[i] if i < len(categories) else "") + "</span>"
            for i in range(len(values))
        )
        return "".join(parts) + '<div class="legend">' + legend + "</div>"

    peak = max([max(s["values"]) for s in series] + [0.0])
    trough = min([min(s["values"]) for s in series] + [0.0])
    span = (peak - trough) or 1.0

    def y_of(value):
        return pad_t + plot_h - ((value - trough) / span) * plot_h

    for step in range(5):
        y = pad_t + (plot_h / 4.0) * step
        parts.append('<line x1="' + str(pad_l) + '" y1="' + str(round(y, 2)) + '" x2="' +
                     str(pad_l + plot_w) + '" y2="' + str(round(y, 2)) +
                     '" stroke="var(--hairline)" stroke-width="1"/>')
        label = trough + span * (1 - step / 4.0)
        parts.append('<text x="' + str(pad_l - 8) + '" y="' + str(round(y + 4, 2)) +
                     '" class="axis" text-anchor="end">' + esc(round(label, 2)) + "</text>")

    slot = plot_w / float(max(len(categories), 1))
    if spec["type"] in ("line", "area"):
        for i, entry in enumerate(series):
            points = []
            for j, value in enumerate(entry["values"]):
                x = pad_l + slot * (j + 0.5)
                points.append(str(round(x, 2)) + "," + str(round(y_of(value), 2)))
            parts.append('<polyline points="' + " ".join(points) + '" fill="none" stroke="#' +
                         PALETTE[i % len(PALETTE)] + '" stroke-width="3" stroke-linejoin="round"/>')
            for point in points:
                x, y = point.split(",")
                parts.append('<circle cx="' + x + '" cy="' + y + '" r="4" fill="#' +
                             PALETTE[i % len(PALETTE)] + '"/>')
    else:
        group = len(series)
        bar_w = max((slot * 0.68) / group, 3.0)
        for i, entry in enumerate(series):
            for j, value in enumerate(entry["values"]):
                x = pad_l + slot * j + (slot - bar_w * group) / 2.0 + bar_w * i
                y = y_of(max(value, 0.0))
                bar_h = abs(y_of(value) - y_of(0.0))
                parts.append('<rect x="' + str(round(x, 2)) + '" y="' + str(round(y, 2)) +
                             '" width="' + str(round(bar_w, 2)) + '" height="' + str(round(max(bar_h, 1.0), 2)) +
                             '" rx="3" fill="#' + PALETTE[i % len(PALETTE)] + '"/>')

    for j, label in enumerate(categories):
        x = pad_l + slot * (j + 0.5)
        parts.append('<text x="' + str(round(x, 2)) + '" y="' + str(height - 16) +
                     '" class="axis" text-anchor="middle">' + esc(label)[:14] + "</text>")
    parts.append("</svg>")

    legend = ""
    if len(series) > 1:
        legend = '<div class="legend">' + "".join(
            '<span class="key"><i style="background:#' + PALETTE[i % len(PALETTE)] + '"></i>' +
            esc(entry["name"]) + "</span>"
            for i, entry in enumerate(series)
        ) + "</div>"
    return "".join(parts) + legend


def table_html(spec):
    head = "".join("<th>" + esc(h) + "</th>" for h in spec["headers"])
    body = "".join(
        "<tr>" + "".join("<td>" + esc(c) + "</td>" for c in row) + "</tr>"
        for row in spec["rows"]
    )
    return '<table class="grid"><thead><tr>' + head + "</tr></thead><tbody>" + body + "</tbody></table>"


def slide_html(name, item, index, total):
    body = ""
    classes = ["slide", "layout-" + name]
    title = esc(item.get("title") or "")
    eyebrow = str(item.get("eyebrow") or "").strip()
    head = ""
    if name not in ("cover", "section", "closing", "quote"):
        head = ('<header>' +
                ('<p class="eyebrow">' + esc(eyebrow.upper()) + "</p>" if eyebrow else "") +
                "<h2>" + title + "</h2><span class=\"rule\"></span></header>")

    if name == "cover":
        meta = [str(v) for v in (item.get("meta") or []) if str(v).strip()]
        if not meta:
            meta = [v for v in (DECK_AUTHOR, str(item.get("date") or PAYLOAD.get("date") or "")) if v]
        body = ('<div class="cover-body">' +
                ('<p class="eyebrow accent">' + esc(eyebrow.upper()) + "</p>" if eyebrow else "") +
                '<span class="rule wide"></span>' +
                "<h1>" + esc(item.get("title") or DECK_TITLE) + "</h1>" +
                ('<p class="sub">' + rich_html(item.get("subtitle") or DECK_SUBTITLE) + "</p>"
                 if (item.get("subtitle") or DECK_SUBTITLE) else "") +
                ('<p class="meta">' + esc("   \u00b7   ".join(meta)) + "</p>" if meta else "") +
                "</div>")
    elif name == "section":
        body = ('<div class="section-body">' +
                ('<span class="section-num">' + esc(item.get("number") or "") + "</span>"
                 if item.get("number") else "") +
                "<h1>" + title + "</h1>" +
                ('<p class="sub">' + rich_html(item.get("body") or item.get("subtitle") or "") + "</p>"
                 if (item.get("body") or item.get("subtitle")) else "") +
                "</div>")
    elif name == "closing":
        entries = normalize_bullets(item.get("bullets"))
        lines = [text for text, _ in entries]
        if item.get("body") or item.get("subtitle"):
            lines.insert(0, str(item.get("body") or item.get("subtitle")))
        body = ('<div class="cover-body">' +
                "<h1>" + esc(item.get("title") or "Thank you") + "</h1>" +
                '<span class="rule wide"></span>' +
                "".join('<p class="sub small">' + rich_html(line) + "</p>" for line in lines[:4]) +
                "</div>")
    elif name == "quote":
        body = ('<div class="quote-body"><span class="mark">\u201c</span>' +
                "<blockquote>" + rich_html(item.get("quote") or item.get("body") or item.get("title") or "") + "</blockquote>" +
                ('<p class="attrib">' + esc(item.get("attribution") or item.get("author") or "") + "</p>"
                 if (item.get("attribution") or item.get("author")) else "") +
                "</div>")
    elif name == "columns":
        columns = [c for c in (item.get("columns") or []) if isinstance(c, dict)][:3]
        cards = ""
        for i, column in enumerate(columns):
            cards += ('<div class="card" style="--card-accent:#' + PALETTE[i % len(PALETTE)] + '">' +
                      ("<h3>" + esc(column.get("heading") or column.get("title") or "") + "</h3>"
                       if (column.get("heading") or column.get("title")) else "") +
                      (bullets_html(normalize_bullets(column.get("bullets") or column.get("points")))
                       or ("<p>" + rich_html(column.get("body") or "") + "</p>" if column.get("body") else "")) +
                      "</div>")
        body = head + '<div class="cards cols-' + str(max(len(columns), 1)) + '">' + cards + "</div>"
    elif name == "stats":
        stats = [s for s in (item.get("stats") or []) if isinstance(s, dict)][:4]
        cards = ""
        for i, stat in enumerate(stats):
            cards += ('<div class="stat" style="--card-accent:#' + PALETTE[i % len(PALETTE)] + '">' +
                      '<span class="value">' + esc(stat.get("value") or stat.get("number") or "") + "</span>" +
                      ('<span class="label">' + esc(stat.get("label") or "") + "</span>" if stat.get("label") else "") +
                      ('<span class="caption">' + esc(stat.get("caption") or stat.get("detail") or "") + "</span>"
                       if (stat.get("caption") or stat.get("detail")) else "") +
                      "</div>")
        body = (head + '<div class="stats cols-' + str(max(len(stats), 1)) + '">' + cards + "</div>" +
                bullets_html(normalize_bullets(item.get("bullets"))[:3]))
    elif name == "chart":
        spec = normalize_chart(item.get("chart"))
        entries = normalize_bullets(item.get("bullets"))
        chart_markup = '<div class="chart-wrap">' + (svg_chart(spec) if spec else "") + "</div>"
        if entries:
            body = (head + '<div class="split"><div class="split-main">' + chart_markup +
                    '</div><aside class="panel">' + bullets_html(entries) + "</aside></div>")
        else:
            body = head + chart_markup
    elif name == "table":
        spec = normalize_table(item.get("table"))
        body = head + (table_html(spec) if spec else "")
    elif name == "image":
        path = resolve_image(item.get("image") or item.get("image_path"))
        uri = data_uri(path) if path else None
        entries = normalize_bullets(item.get("bullets"))
        figure = ('<figure><img src="' + uri + '" alt="' + esc(item.get("caption") or title) + '"/>' +
                  ("<figcaption>" + esc(item.get("caption")) + "</figcaption>" if item.get("caption") else "") +
                  "</figure>") if uri else '<div class="panel muted">Image not embedded in preview</div>'
        if entries:
            body = (head + '<div class="split"><div class="split-main">' + figure +
                    '</div><aside class="panel">' + bullets_html(entries) + "</aside></div>")
        else:
            body = head + figure
    else:
        entries = normalize_bullets(item.get("bullets"))
        intro = ('<p class="intro">' + rich_html(item.get("body")) + "</p>") if item.get("body") else ""
        body = head + intro + bullets_html(entries)

    footer = ""
    if name not in ("cover",):
        footer = ('<footer><span>' + esc(FOOTER_LABEL) + "</span><span>" +
                  str(index) + " / " + str(total) + "</span></footer>")
    return ('<div class="frame"><section class="' + " ".join(classes) + '">' +
            body + footer + "</section></div>")


if HTML_PATH:
    # Built by concatenation on purpose: the stylesheet below is full of literal
    # percent signs, so %-formatting or str.format would need escaping everywhere.
    def var(name, key, fallback):
        return "--" + name + ":#" + str(THEME.get(key) or fallback).lstrip("#") + ";"

    root_css = (
        ":root{" +
        var("bg", "background", "FFFFFF") +
        var("surface", "surface", "F4F6FE") +
        var("surface-alt", "surface_alt", "EAEEFC") +
        var("ink", "ink", "0F172A") +
        var("muted", "muted", "64748B") +
        var("hairline", "hairline", "DDE3F0") +
        var("accent", "accent", "4F46E5") +
        var("accent-soft", "accent_soft", "E5E8FD") +
        var("cover-bg", "cover_background", "1E1B4B") +
        var("cover-ink", "cover_ink", "FFFFFF") +
        var("cover-muted", "cover_muted", "C7D2FE") +
        var("cover-accent", "cover_accent", "A5B4FC") +
        "--head:" + (FONTS.get("css_heading") or "sans-serif") + ";" +
        "--body:" + (FONTS.get("css_body") or "sans-serif") + ";" +
        "--mono:" + (FONTS.get("css_mono") or "monospace") + ";" +
        "}"
    )

    css = root_css + """
*{box-sizing:border-box}
html,body{margin:0;height:100%;background:#0b1220}
body{font-family:var(--body);-webkit-font-smoothing:antialiased}
.deck{height:100%;overflow-y:auto;scroll-snap-type:y mandatory;padding:16px;display:flex;
      flex-direction:column;gap:16px;align-items:center}
/* .frame is the query container so that cqw units on .slide itself resolve
   against the slide width. An element can never query its own size. */
.frame{width:100%;max-width:1120px;flex:0 0 auto;scroll-snap-align:center;container-type:inline-size}
.slide{position:relative;aspect-ratio:16/9;background:var(--bg);color:var(--ink);border-radius:12px;
       overflow:hidden;padding:4.4cqw 5.2cqw 7.2cqw;display:flex;flex-direction:column;
       justify-content:flex-start;box-shadow:0 18px 50px rgba(0,0,0,.45)}
.slide>*{position:relative;z-index:1}
header{margin-bottom:2.2cqw}
h1,h2,h3{font-family:var(--head);margin:0;letter-spacing:-.015em;line-height:1.1}
h2{font-size:3.1cqw;font-weight:700}
h1{font-size:4.6cqw;font-weight:700}
h3{font-size:1.9cqw;font-weight:700;margin-bottom:1.1cqw}
.eyebrow{margin:0 0 .7cqw;font-size:1.15cqw;font-weight:700;letter-spacing:.14em;color:var(--accent)}
.rule{display:block;width:9cqw;height:.42cqw;background:var(--accent);border-radius:99px;margin-top:1.4cqw}
.rule.wide{width:14cqw;margin:0 0 2.2cqw}
.bullets{list-style:none;margin:0;padding:0;font-size:1.72cqw;line-height:1.5}
.bullets li{position:relative;padding-left:2.2cqw;margin:.85cqw 0;color:var(--ink)}
.bullets li::before{content:"";position:absolute;left:0;top:.55em;width:.78cqw;height:.78cqw;
                    background:var(--accent);border-radius:2px}
.bullets li.lvl1{margin-left:2.2cqw;font-size:1.5cqw;color:var(--muted)}
.bullets li.lvl1::before{background:var(--muted);height:.14cqw;top:.72em;width:1cqw;border-radius:99px}
.bullets li.lvl2{margin-left:4.4cqw;font-size:1.4cqw;color:var(--muted)}
.intro{font-size:1.6cqw;color:var(--muted);margin:0 0 1.6cqw;line-height:1.5}
code{font-family:var(--mono);background:var(--accent-soft);color:var(--accent);
     padding:.1em .35em;border-radius:4px;font-size:.92em}
.layout-cover,.layout-section,.layout-closing{background:var(--cover-bg);color:var(--cover-ink);
       justify-content:center}
.layout-cover h1,.layout-section h1,.layout-closing h1{color:var(--cover-ink)}
.layout-cover .rule,.layout-closing .rule{background:var(--cover-accent)}
.layout-cover .eyebrow.accent{color:var(--cover-accent)}
.layout-cover .sub,.layout-section .sub,.layout-closing .sub{font-size:2cqw;color:var(--cover-muted);
       margin:1.6cqw 0 0;max-width:72%;line-height:1.45}
.layout-closing .sub.small{font-size:1.6cqw;margin:.7cqw 0 0}
.layout-closing .rule.wide{margin:2.2cqw 0 1.6cqw}
.layout-cover .meta{margin-top:3cqw;font-size:1.2cqw;color:var(--cover-muted);letter-spacing:.05em}
.layout-cover::after{content:"";position:absolute;right:-8cqw;top:-14cqw;width:46cqw;aspect-ratio:1;
       border-radius:50%;background:var(--cover-accent);opacity:.12;z-index:0}
.layout-section{border-left:.7cqw solid var(--cover-accent)}
.section-num{display:block;font-family:var(--head);font-size:6.4cqw;font-weight:700;
       color:var(--cover-accent);line-height:1;margin-bottom:1.2cqw}
.layout-quote{background:var(--surface);justify-content:center;border-left:.7cqw solid var(--accent)}
.quote-body{position:relative}
.mark{position:absolute;top:-6.4cqw;left:-.8cqw;font-family:var(--head);font-size:11cqw;
      color:var(--accent);opacity:.22;line-height:1}
blockquote{margin:0;font-family:var(--head);font-style:italic;font-size:2.9cqw;line-height:1.32;
      color:var(--ink);max-width:82%}
.attrib{margin:2.4cqw 0 0;font-size:1.35cqw;color:var(--muted);padding-top:1.2cqw;
      border-top:.28cqw solid var(--accent);display:inline-block}
.cards,.stats{display:grid;gap:1.6cqw;flex:1;min-height:0}
.cols-1{grid-template-columns:1fr}
.cols-2{grid-template-columns:1fr 1fr}
.cols-3{grid-template-columns:repeat(3,1fr)}
.cols-4{grid-template-columns:repeat(4,1fr)}
.card{background:var(--surface);border-radius:.9cqw;padding:2cqw;border-top:.42cqw solid var(--card-accent);
      overflow:hidden}
.card p{font-size:1.5cqw;color:var(--muted);line-height:1.5;margin:0}
.card .bullets{font-size:1.42cqw}
.stat{background:var(--surface);border-radius:.9cqw;padding:2cqw;border-left:.42cqw solid var(--card-accent);
      display:flex;flex-direction:column;justify-content:center;gap:.5cqw}
.stat .value{font-family:var(--head);font-size:3.9cqw;font-weight:700;color:var(--card-accent);line-height:1}
.stat .label{font-size:1.4cqw;font-weight:700;color:var(--ink)}
.stat .caption{font-size:1.12cqw;color:var(--muted);line-height:1.35}
.split{display:grid;grid-template-columns:1.62fr 1fr;gap:1.8cqw;flex:1;min-height:0}
.split-main{min-width:0;display:flex;align-items:center;justify-content:center}
.panel{background:var(--surface);border-radius:.9cqw;padding:1.8cqw;overflow:hidden}
.panel.muted{color:var(--muted);display:flex;align-items:center;justify-content:center;font-size:1.4cqw}
.chart-wrap{flex:1;min-height:0;display:flex;flex-direction:column;align-items:center;justify-content:center}
svg.chart{width:100%;height:auto;max-height:100%}
.axis{font-family:var(--body);font-size:11px;fill:var(--muted)}
.legend{display:flex;flex-wrap:wrap;gap:1.4cqw;justify-content:center;margin-top:.8cqw;
      font-size:1.15cqw;color:var(--muted)}
.key{display:inline-flex;align-items:center;gap:.5cqw}
.key i{width:.9cqw;height:.9cqw;border-radius:2px;display:inline-block}
table.grid{width:100%;border-collapse:collapse;font-size:1.32cqw}
table.grid th{background:var(--accent);color:#fff;text-align:left;padding:.9cqw 1.1cqw;font-weight:700}
table.grid td{padding:.8cqw 1.1cqw;border-bottom:1px solid var(--hairline);color:var(--ink)}
table.grid tbody tr:nth-child(even){background:var(--surface)}
figure{margin:0;flex:1;min-height:0;display:flex;flex-direction:column;align-items:center;
      justify-content:center;gap:1cqw}
figure img{max-width:100%;max-height:100%;object-fit:contain;border-radius:.7cqw}
figcaption{font-size:1.2cqw;color:var(--muted)}
.slide>footer{position:absolute;left:5.2cqw;right:5.2cqw;bottom:2.8cqw;display:flex;
      justify-content:space-between;font-size:1.05cqw;color:var(--muted);
      border-top:1px solid var(--hairline);padding-top:.9cqw}
@media print{.deck{display:block;padding:0}.frame{max-width:none;page-break-after:always}
      .slide{box-shadow:none;border-radius:0}}
"""

    script = (
        "<script>document.addEventListener('keydown',function(e){"
        "var slides=[].slice.call(document.querySelectorAll('.frame'));"
        "if(!slides.length)return;"
        "var deck=document.querySelector('.deck');"
        "var mid=deck.scrollTop+deck.clientHeight/2;var cur=0;"
        "slides.forEach(function(s,i){if(s.offsetTop<=mid)cur=i;});"
        "var next=cur;"
        "if(e.key==='ArrowDown'||e.key==='ArrowRight'||e.key===' ')next=Math.min(cur+1,slides.length-1);"
        "else if(e.key==='ArrowUp'||e.key==='ArrowLeft')next=Math.max(cur-1,0);else return;"
        "e.preventDefault();slides[next].scrollIntoView({behavior:'smooth',block:'center'});"
        "});</script>"
    )

    sections = "".join(
        slide_html(name, item, i + 1, total) for i, (name, item) in enumerate(prepared)
    )
    document = (
        '<!DOCTYPE html><html lang="en"><head><meta charset="utf-8">'
        '<meta name="viewport" content="width=device-width, initial-scale=1">'
        "<title>" + esc(DECK_TITLE) + "</title><style>" + css + "</style></head>"
        '<body><div class="deck">' + sections + "</div>" + script + "</body></html>"
    )
    os.makedirs(os.path.dirname(HTML_PATH), exist_ok=True)
    with open(HTML_PATH, "w", encoding="utf-8") as handle:
        handle.write(document)

result = {
    "status": "success",
    "path": OUT_PATH,
    "size": os.path.getsize(OUT_PATH),
    "slide_count": total,
    "layouts": [name for name, _ in prepared],
    "theme": THEME.get("id") or "aurora",
}
if HTML_PATH and os.path.exists(HTML_PATH):
    result["html_path"] = HTML_PATH
    result["html_size"] = os.path.getsize(HTML_PATH)
print(json.dumps(result))
'''
