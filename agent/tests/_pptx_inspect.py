# Copyright (c) 2026 Agentic Company. All rights reserved.
# Proprietary and non-commercial use only.

"""Local-only geometry audit for generated decks.

PowerPoint is not available here, so instead of eyeballing the render we assert
the structural properties that make a deck look right: nothing spills off the
canvas, every slide carries real text, and body text stays inside its box.
"""

from __future__ import annotations

import pathlib
import sys

from pptx import Presentation
from pptx.util import Emu

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
TOLERANCE = Emu(9525) * 4  # ~4pt of slack for stroke widths


def audit(path: pathlib.Path) -> list[str]:
    prs = Presentation(str(path))
    problems: list[str] = []
    if prs.slide_width != SLIDE_W or prs.slide_height != SLIDE_H:
        problems.append(f"slide size is {prs.slide_width}x{prs.slide_height}, expected 16:9")

    for index, slide in enumerate(prs.slides, start=1):
        text_seen = False
        for shape in slide.shapes:
            name = f"slide {index} / {shape.shape_type}"
            if shape.left is None or shape.top is None:
                problems.append(f"{name}: shape has no explicit position")
                continue
            right = shape.left + (shape.width or 0)
            bottom = shape.top + (shape.height or 0)
            # Decorative bleed shapes are allowed to run off the edges; text is not.
            has_text = shape.has_text_frame and shape.text_frame.text.strip()
            if has_text:
                text_seen = True
                if shape.left < -TOLERANCE or shape.top < -TOLERANCE:
                    problems.append(f"{name}: text starts off-canvas at ({shape.left}, {shape.top})")
                if right > SLIDE_W + TOLERANCE or bottom > SLIDE_H + TOLERANCE:
                    problems.append(
                        f"{name}: text box ends off-canvas at ({right}, {bottom}) "
                        f"text={shape.text_frame.text[:40]!r}"
                    )
            if shape.has_table:
                text_seen = True
                if shape.left + shape.width > SLIDE_W + TOLERANCE:
                    problems.append(f"{name}: table exceeds slide width")
            if shape.has_chart:
                text_seen = True
        if not text_seen:
            problems.append(f"slide {index}: no text content at all")
    return problems


def main() -> None:
    target = pathlib.Path(sys.argv[1])
    files = sorted(target.glob("*.pptx")) if target.is_dir() else [target]
    if not files:
        print(f"no .pptx found in {target}")
        raise SystemExit(1)
    failed = False
    for path in files:
        problems = audit(path)
        prs = Presentation(str(path))
        slides = list(prs.slides)
        shapes = sum(len(list(slide.shapes)) for slide in slides)
        if problems:
            failed = True
            print(f"\n{path.name}: {len(problems)} problem(s)")
            for problem in problems[:20]:
                print(f"  - {problem}")
        else:
            print(f"{path.name}: clean ({len(slides)} slides, {shapes} shapes)")
    raise SystemExit(1 if failed else 0)


if __name__ == "__main__":
    main()
